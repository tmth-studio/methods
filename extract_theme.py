#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
THEME EXTRACTOR — pull a client brand out of any .pptx, Copilot-executable.

Companion to the pyramid and proposition generators in this repository. Every
generator holds its brand as data in a THEME dict (neutral TMTH palette by
default). This script reads an attached client .pptx and prints a
ready-to-paste THEME dict, so client brand arrives at runtime and never lives
in public code.

HOW TO USE (in Copilot, Claude, or any assistant with a code interpreter):
  1. Attach this file AND any existing client .pptx (a real deck the client
     produced — a template, a recent board pack, anything on-brand).
  2. Say: "Run this script." It finds the attached .pptx automatically
     (or set INPUT_FILE below to a specific filename).
  3. It prints: (a) the raw theme — the file's colour scheme (dk1/dk2,
     lt1/lt2, accent1–6) and major/minor fonts; (b) the colours actually
     observed on the first slides; (c) a SUGGESTED THEME dict mapped onto
     the generators' semantic roles. Paste the dict into a generator's DATA
     section, then override any role by hand — the suggestion is a starting
     point, and the raw theme is printed precisely so you can overrule it.

WHAT THE ROLES MEAN (identical across all generators):
  dominant       lead brand colour — covers, dividers, emphasis, table headers
  accent         secondary brand colour — titles, chart lines
  highlight      the one-highlight colour — eyebrows, numbers, the ONE
                 highlighted element per exhibit
  card_fill      card / panel background tint (must stay light)
  ink            near-black headings ink
  ink_soft       body text ink
  muted          secondary text, sources, footers
  chart_compare  non-highlighted chart bars (must stay quiet)
  font_heading / font_body

DEGRADED MODE: if the .pptx has no readable theme part, or its theme is the
stock Office default (common in generated decks), the script says so and
builds the suggestion from the colours observed on the slides instead.

VERSION NOTES
  v1.0 — 20 Aug 2026: first version — theme-part extraction (srgbClr and
         sysClr), stock-Office-theme detection, observed-colour fallback
         from slide fills and text runs, suggested role mapping.

Requires: python-pptx.  Output: printed report (nothing written to disk).
"""

# ============================================================================
# DATA — EDIT THIS SECTION ONLY
# ============================================================================

# Leave empty to auto-detect: the script takes the most recently modified
# .pptx in the working directory (excluding the generators' own outputs).
INPUT_FILE = ""

# How many slides to scan for observed colours and fonts.
SLIDES_TO_SCAN = 6

# ============================================================================
# BUILD — DO NOT EDIT BELOW THIS LINE
# ============================================================================
import collections
import glob
import os
import sys
from lxml import etree

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Stock Office themes (2007-era and 2013+ defaults); if the file carries one,
# the theme part tells us nothing about the client's brand.
STOCK_OFFICE_ACCENTS = {"4472C4", "5B9BD5", "70AD47", "ED7D31", "FFC000",
                        "A5A5A5", "44546A", "E7E6E6",
                        "4F81BD", "C0504D", "9BBB59", "8064A2", "4BACC6",
                        "F79646"}
OWN_OUTPUTS = {"presented-deck.pptx", "standalone-deck.pptx",
               "proposition-backlog-template-branded.pptx",
               "partner-decision-template-branded.pptx"}


def fail(msg):
    print("=" * 64)
    print("SELF-CHECK — extract_theme.py")
    print("  [FAIL] %s" % msg)
    print("=" * 64)
    raise SystemExit("EXTRACTION FAILED — nothing printed that you should "
                     "paste. Fix the input and re-run.")


def pick_input():
    if INPUT_FILE:
        if not os.path.exists(INPUT_FILE):
            fail("INPUT_FILE %r does not exist in the working directory"
                 % INPUT_FILE)
        return INPUT_FILE
    cands = [p for p in glob.glob("*.pptx")
             if os.path.basename(p) not in OWN_OUTPUTS
             and not os.path.basename(p).startswith("~$")]
    if not cands:
        cands = [p for p in glob.glob("*.pptx")
                 if not os.path.basename(p).startswith("~$")]
    if not cands:
        fail("no .pptx found — attach the client deck alongside this script,"
             " or set INPUT_FILE")
    return max(cands, key=os.path.getmtime)


def read_theme_part(prs):
    """Return ({slot: hex}, {'major': font, 'minor': font}) or (None, None)."""
    try:
        from pptx.opc.constants import RELATIONSHIP_TYPE as RT
        part = prs.slide_masters[0].part.part_related_by(RT.THEME)
        root = etree.fromstring(part.blob)
    except Exception as e:
        print("  note: no readable theme part (%s)" % e.__class__.__name__)
        return None, None
    colors = {}
    scheme = root.find(".//{%s}clrScheme" % A_NS)
    if scheme is None:
        return None, None
    for child in scheme:
        slot = etree.QName(child).localname
        srgb = child.find("{%s}srgbClr" % A_NS)
        sysc = child.find("{%s}sysClr" % A_NS)
        if srgb is not None:
            colors[slot] = srgb.get("val", "").upper()
        elif sysc is not None:
            colors[slot] = (sysc.get("lastClr") or
                            {"windowText": "000000",
                             "window": "FFFFFF"}.get(sysc.get("val", ""),
                                                     "")).upper()
    fonts = {}
    for tag, key in (("majorFont", "major"), ("minorFont", "minor")):
        el = root.find(".//{%s}%s/{%s}latin" % (A_NS, tag, A_NS))
        if el is not None and el.get("typeface"):
            fonts[key] = el.get("typeface")
    return (colors or None), (fonts or None)


def scan_slides(prs, n):
    """Observed solid-fill colours, text colours and run fonts, by count."""
    fills = collections.Counter()
    text_colors = collections.Counter()
    fonts = collections.Counter()
    font_max_size = {}

    def walk(shapes):
        for shp in shapes:
            if shp.shape_type == 6:  # group
                try:
                    walk(shp.shapes)
                except Exception:
                    pass
                continue
            try:
                if (shp.fill.type is not None and shp.fill.type == 1
                        and shp.fill.fore_color.type == 1):
                    fills[str(shp.fill.fore_color.rgb)] += 1
            except Exception:
                pass
            if getattr(shp, "has_text_frame", False):
                for p in shp.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.name:
                            fonts[r.font.name] += 1
                            if r.font.size is not None:
                                font_max_size[r.font.name] = max(
                                    font_max_size.get(r.font.name, 0),
                                    r.font.size.pt)
                        try:
                            if r.font.color and r.font.color.type == 1:
                                text_colors[str(r.font.color.rgb)] += 1
                        except Exception:
                            pass

    for sl in list(prs.slides)[:n]:
        walk(sl.shapes)
    return fills, text_colors, fonts, font_max_size


def luminance(h):
    r, g, b = (int(h[i:i + 2], 16) / 255.0 for i in (0, 2, 4))
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def saturation(h):
    r, g, b = (int(h[i:i + 2], 16) / 255.0 for i in (0, 2, 4))
    mx, mn = max(r, g, b), min(r, g, b)
    return 0.0 if mx == 0 else (mx - mn) / mx


def blend(base, mix, f):
    b = [int(base[i:i + 2], 16) for i in (0, 2, 4)]
    m = [int(mix[i:i + 2], 16) for i in (0, 2, 4)]
    return "".join("%02X" % round(b[i] * (1 - f) + m[i] * f)
                   for i in range(3))


def suggest_from_theme(colors):
    acc = [colors.get("accent%d" % i) for i in range(1, 7)]
    acc = [a for a in acc if a]
    s = {}
    s["dominant"] = colors.get("dk2") or (acc[0] if acc else "2F4858")
    # prefer the strongest accents for the visible roles
    s["accent"] = acc[0] if acc else "4E7C90"
    s["highlight"] = acc[1] if len(acc) > 1 else "C9A227"
    lt2 = colors.get("lt2", "")
    s["card_fill"] = lt2 if lt2 and luminance(lt2) > 0.80 else \
        blend(s["dominant"], "FFFFFF", 0.90)
    s["ink"] = colors.get("dk1") or "1F2A31"
    return s


def suggest_from_observed(fills, text_colors):
    # fills carry double weight — they are deliberate brand surfaces;
    # text colours widen the candidate pool (titles, eyebrows, numerals).
    merged = collections.Counter()
    for h, c in fills.items():
        merged[h] += 2 * c
    for h, c in text_colors.items():
        merged[h] += c
    pal = [(h, c) for h, c in merged.most_common(24)
           if h not in ("FFFFFF", "000000")]
    dark = [h for h, _ in pal if luminance(h) < 0.35]
    mid = [h for h, _ in pal if 0.35 <= luminance(h) < 0.75]
    light = [h for h, _ in pal if luminance(h) >= 0.75]
    s = {}
    s["dominant"] = dark[0] if dark else "2F4858"
    # the most saturated non-dominant candidate reads as the highlight;
    # the accent is the next most frequent mid/dark tone after those two
    by_sat = sorted(mid + dark, key=saturation, reverse=True)
    s["highlight"] = next((h for h in by_sat if h != s["dominant"]),
                          "C9A227")
    s["accent"] = next((h for h in mid + dark
                        if h not in (s["dominant"], s["highlight"])),
                       "4E7C90" if s["highlight"] != "4E7C90"
                       else "6B7A85")
    s["card_fill"] = light[0] if light else "EEF3F6"
    inks = [h for h, _ in text_colors.most_common(12)
            if luminance(h) < 0.25]
    s["ink"] = inks[0] if inks else "1F2A31"
    return s


def finish(s):
    """Fill the derived roles from what was chosen."""
    s.setdefault("ink", "1F2A31")
    s["ink_soft"] = blend(s["ink"], s["dominant"], 0.35)
    s["muted"] = blend(s["ink"], "FFFFFF", 0.45)
    s["chart_compare"] = blend(s["dominant"], "FFFFFF", 0.62)
    return s


def main():
    path = pick_input()
    print("reading:", path)
    try:
        from pptx import Presentation
        prs = Presentation(path)
    except Exception as e:
        fail("could not open %r as a .pptx (%s)" % (path, e))

    theme_colors, theme_fonts = read_theme_part(prs)
    fills, text_colors, run_fonts, font_max_size = scan_slides(
        prs, SLIDES_TO_SCAN)

    stock = bool(theme_colors) and (
        theme_colors.get("accent1", "") in STOCK_OFFICE_ACCENTS
        and theme_colors.get("accent2", "") in STOCK_OFFICE_ACCENTS)

    print()
    print("RAW THEME PART" + ("  (none readable)" if not theme_colors
                              else ""))
    if theme_colors:
        for slot in ("dk1", "lt1", "dk2", "lt2", "accent1", "accent2",
                     "accent3", "accent4", "accent5", "accent6", "hlink",
                     "folHlink"):
            if slot in theme_colors:
                print("  %-9s %s" % (slot, theme_colors[slot]))
        if theme_fonts:
            print("  major font: %s   minor font: %s"
                  % (theme_fonts.get("major", "—"),
                     theme_fonts.get("minor", "—")))
        if stock:
            print("  note: this is the STOCK Office theme — it says nothing "
                  "about the client's brand; using the colours observed on "
                  "the slides instead.")
    print()
    print("OBSERVED ON THE FIRST %d SLIDES (colour: count)"
          % min(SLIDES_TO_SCAN, len(prs.slides._sldIdLst)))
    for h, c in fills.most_common(12):
        print("  fill %s: %d" % (h, c))
    for h, c in text_colors.most_common(8):
        print("  text %s: %d" % (h, c))
    for f, c in run_fonts.most_common(6):
        print("  font %s: %d" % (f, c))

    use_theme_part = bool(theme_colors) and not stock
    if use_theme_part:
        s = suggest_from_theme(theme_colors)
        basis = "the file's theme part"
    else:
        if not fills and not text_colors:
            fail("no theme part and no observable colours on the slides — "
                 "attach a deck with real content")
        s = suggest_from_observed(fills, text_colors)
        basis = ("the colours observed on the slides (no usable theme "
                 "part)")
    s = finish(s)

    if use_theme_part and theme_fonts:
        head_font = theme_fonts.get("major") or "Cambria"
        body_font = theme_fonts.get("minor") or head_font
    else:
        # the font carrying the largest observed size reads as the heading
        # face; the most frequent other font is the body face
        common = [f for f, _ in run_fonts.most_common(4)]
        if font_max_size:
            head_font = max(font_max_size, key=font_max_size.get)
        else:
            head_font = common[0] if common else "Cambria"
        body_font = next((f for f in common if f != head_font),
                         head_font)

    # sanity: every suggested colour must be a valid 6-char hex
    bad = [k for k, v in s.items() if not (isinstance(v, str) and len(v) == 6
           and all(c in "0123456789ABCDEFabcdef" for c in v))]
    if bad:
        fail("internal: malformed suggestion for %s — report this deck"
             % ", ".join(bad))

    print()
    print("SUGGESTED THEME — built from %s." % basis)
    print("Paste into a generator's DATA section; override any role by "
          "hand from the raw values above.")
    print()
    print("THEME = {")
    for k in ("dominant", "accent", "highlight", "card_fill", "ink",
              "ink_soft", "muted", "chart_compare"):
        print('    "%s":%s"%s",' % (k, " " * (14 - len(k)), s[k]))
    print('    "font_heading":  "%s",' % head_font)
    print('    "font_body":     "%s",' % body_font)
    print("}")
    print()
    print("CHECK BEFORE USING: card_fill must stay light (it sits behind "
          "text), chart_compare must stay quiet (it is the NON-highlighted "
          "bar), and highlight must contrast with both. The generators "
          "validate hex shape, not taste — look at one rendered page "
          "before shipping.")


main()
