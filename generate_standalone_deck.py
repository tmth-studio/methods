#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
STANDALONE DECK GENERATOR — self-sufficient Minto slides that read unaided.

Produces a .pptx that travels WITHOUT a presenter: an executive summary slide
carries the full SCQA in prose up front, every slide body proves its title on
the page, and every evidence slide cites its source on the slide. For a deck
a speaker will present live, use generate_presentation_deck.py — it takes the
identical DATA.

HOW TO USE (in Copilot, Claude, or any assistant with a code interpreter):
  1. Attach this file. Say: "Run this script and give me the .pptx it creates."
     It ships with fictional demonstration content (Meridian Office Services).
  2. To build your own deck: the assistant follows the STAGED PROMPTS below
     and edits ONLY the DATA section. Everything under BUILD — DO NOT EDIT
     stays. Re-run after every edit; the SELF-CHECK refuses to write a file
     while any rule fails.

THE SHARED SCHEMA: the PYRAMID data below is identical to the DATA sections of
generate_pyramid_memo.py and generate_presentation_deck.py. Fill it once,
paste it between the DATA markers of all three, and one argument drives the
memo and both decks. (Each script sets its own output filename in the build
section, so a pasted DATA block never renames the artefact.) Generator-
specific registers sit at the END of each DATA section (the decks:
EXHIBIT_REGISTER + MAX_VINTAGE_MONTHS; the memo: EVIDENCE_REGISTER) — when
pasting the shared block between scripts, keep each script's own registers
in place.

STAGED PROMPTS — the assistant follows these in order and STOPS at each gate.

  STAGE 0 — APPLY THE CLIENT'S BRAND (optional). All brand values live in
  the THEME dict at the top of the DATA section; the shipped default is the
  neutral TMTH palette. Two routes to a client brand: (i) paste the client's
  palette hexes and font names into THEME; (ii) attach any existing client
  .pptx and run extract_theme.py (same repository) — it prints a
  ready-to-paste THEME dict with a suggested role mapping. Client brand
  arrives at runtime; never commit a client THEME to the public repository.

  STAGE 1 — FRAME. Ask, one question at a time: who reads this, and where —
  inbox, portal, pre-read? What decision is wanted? Which of the four patterns
  (giving direction / seeking approval / explaining how / choosing)? Confirm
  the deck travels WITHOUT a presenter — if the author presents live, switch
  to the presented-deck generator. GATE: get a yes on the frame.

  STAGE 2 — SCQA AND GOVERNING THOUGHT. Draft Situation, Complication, one
  Question, and the governing thought in three drafts (WHAT → + SO-WHAT →
  25 words or fewer, active, verb first). GATE: show the drafts, get a yes.

  STAGE 3 — PYRAMID. 2–4 key-line arguments as full-sentence assertions; name
  the ordering principle; MECE test both ways; 2–5 supports per argument from
  evidence the user supplies — never invented. GATE: show the pyramid, get a
  yes.

  STAGE 4 — EXHIBIT REUSE. Attach the engagement's prior decks. Catalogue
  the existing exhibits — for each: the MESSAGE it proves (a claim, not a
  topic), its chart form, its source deck, and its data vintage (the date
  of the underlying data, not the deck date). For every key-line point
  needing an exhibit, propose REUSE / ADAPT (say exactly what changes) /
  NEW, and fill the EXHIBIT_REGISTER. Staleness guard: any REUSE or ADAPT
  whose data vintage is older than MAX_VINTAGE_MONTHS (default 6) FAILS the
  build unless flagged vintage_accepted with a stated reason. Only NEW and
  ADAPT exhibits get design effort in Stage 5. GATE: show the register,
  get a yes.

  STAGE 5 — BODIES AND EXHIBITS. For each support, write the on-slide prose:
  two to four sentences (15–55 words) that prove the title to a cold reader
  with nobody in the room. Claim first, then the exhibit: name the comparison
  (component / item / time series / frequency / correlation), pick the chart
  kind — bar / column / line (one or two series) / waterfall (build-ups and
  bridges) / table — supply real data with a source, and give each exhibit
  its register exhibit_id. If a title claims two things moved ("X fell while
  Y rose"), the exhibit must show BOTH series — index them to a common base
  year when their units differ. The three sidebar bullets must add a second
  layer (the so-what, the caveat, the scale comparison) — never restate what
  the exhibit already shows. GATE: show three sample slide bodies, get a yes.

  STAGE 6 — FILL AND RUN. Write the approved content into DATA. Replace every
  square-bracketed placeholder — the build fails while any remain. Run the
  script; hand back the .pptx and the self-check report.

ENGAGEMENT ASSETS: a filled EXHIBIT_REGISTER is client work product. It
lives in the client tenant. Never commit a filled register to the public
repository — the shipped register holds fictional demonstration content only.

VERSION NOTES
  v1.2 — 20 Aug 2026: consulting-density pass. Two-series charts (line /
         column / bar with a "series" list) and a waterfall kind for
         build-ups and bridges; legend appears only on multi-series charts.
         Tighter body geometry (exhibit moved up, dead band removed),
         no-exhibit supports laid out as three parallel cards, page numbers
         on content slides. Geometry guards in the self-check: titles 5–15
         words, prose 15–55 words, exactly 1–3 bullets, why_it_matters
         40 words max (previously long bullets were silently truncated).
         New warning when a support quotes numbers but carries no exhibit.
         Default fonts now Georgia / Arial (travel-safe on client machines).
  v1.1 — 20 Aug 2026: brand moved from code to data (THEME dict, neutral
         TMTH default; validation gate on hex/fonts; dark-slide tints now
         derived from the dominant colour). Exhibit-reuse stage added with
         the EXHIBIT_REGISTER, exhibit_id linking, and a data-vintage
         staleness guard (MAX_VINTAGE_MONTHS, default 6).
  v1.0 — 20 Aug 2026: first published version.

Requires: python-pptx.  Output: standalone-deck.pptx
"""

# ============================================================================
# DATA — EDIT THIS SECTION ONLY  (shared PYRAMID schema — see docstring)
# All demonstration content is fictional (Meridian Office Services).
# ============================================================================

# --- THEME: all brand values live here, as data. Ship = neutral TMTH palette.
# To apply a client's brand: (i) paste their palette hexes + font names below,
# or (ii) attach any existing client .pptx and run extract_theme.py (same
# repository) — it prints a ready-to-paste THEME dict. Hex: 6 chars, no '#',
# no alpha. Never commit a client THEME to the public repository.
THEME = {
    "dominant":      "2F4858",  # lead brand colour — cover, dividers, the answer, table headers
    "accent":        "4E7C90",  # secondary brand colour — slide titles, line charts
    "highlight":     "C9A227",  # the one-highlight colour — eyebrows, numbers, highlighted bars
    "card_fill":     "EEF3F6",  # card / panel background tint
    "ink":           "1F2A31",  # near-black headings ink
    "ink_soft":      "3C4C57",  # body text ink
    "muted":         "6B7A85",  # secondary text, sources, footers
    "chart_compare": "B9C6CD",  # non-highlighted chart bars / comparison elements
    "font_heading":  "Georgia",  # travel-safe serif — installed on Windows and Mac alike
    "font_body":     "Arial",    # travel-safe sans — never substituted on client machines
}

META = {
    "title": "Win renewals on response time, not price",
    "subtitle": "Recommendation — SME contract churn",
    "author_line": "Strategy team",
    "date": "20 Aug 2026",
    "audience": "Managing director and the executive committee",
    "presence": "standalone",        # this generator assumes no presenter
    "archetype": "recommendation",   # used by the memo generator; ignored here
    "footer_tag": "Demonstration content — every name and number is fictional",
}

SCQA = {
    "tone_order": "SCA",  # SCA standard | ASC direct | CSA concerned
    "situation": ("Meridian serves 1,200 SME service contracts across the "
                  "South East, worth £5.1M of annual revenue. Renewal pricing "
                  "has been flat for two years."),
    "complication": ("Churn has risen from 22% to 35% in three years — a "
                     "£1.8M annual revenue impact. The lost-deal evidence "
                     "points away from price: objections fell 30% while churn "
                     "climbed."),
    "question": ("How does Meridian win renewals in a market that now judges "
                 "providers on response time?"),
}

GOVERNING_THOUGHT = ("Reposition Meridian around a response-time guarantee — "
                     "the factor deciding renewals — recovering churn to 25% "
                     "and £1.2M annual revenue within four quarters.")

KEY_LINE_ORDER = "degree"  # "time" | "structure" | "degree"

KEY_LINE = [
    {
        "tag": "DIAGNOSIS",
        "assertion": "Churn is perception-driven, not price-driven",
        "why_it_matters": ("Every remedy on the table prices the wrong "
                           "problem. The evidence separates what customers "
                           "say from what they do — and it clears price as "
                           "the cause."),
        "supports": [
            {
                "assertion": ("Current customers renew at stable rates while "
                              "competitor-bid churn climbs"),
                "prose": ("Existing-contract churn has held between 21% and "
                          "23% for eight quarters. A price problem hits "
                          "paying customers first, and it has not. The rise "
                          "is concentrated in renewals where a competitor "
                          "bid is on the table."),
                "bullets": [
                    "Existing-contract churn steady at 21–23% for eight quarters",
                    "Price pain reaches current payers first — it has not",
                    "The rise sits in renewals facing a competitor bid",
                ],
                "source": "Contract system renewal records, Q3 2024 to Q2 2026, n=1,180",
                "exhibit": None,
                "notes": ("Land the distinction between the two churn "
                          "populations before anyone mentions price. The "
                          "whole argument turns on it."),
            },
            {
                "assertion": ("Price objections fell 30% while churn rose 13 "
                              "points"),
                "prose": ("Recorded price objections fell from 46 to 32 a "
                          "quarter across three years. If price drove the "
                          "losses, objections would rise with churn. They "
                          "moved the opposite way."),
                "bullets": [
                    "A price-driven loss pattern would move these lines together",
                    "The divergence widens every year — this is not noise",
                    "Objections data covers all 1,200 contracts, not a sample",
                ],
                "source": "CRM objection log and renewal records, 2023 to 2026",
                "exhibit": {"exhibit_id": "EX-OBJECTIONS",
                            "kind": "line",
                            "unit": "Indexed, 2023 = 100",
                            "categories": ["2023", "2024", "2025", "2026"],
                            "series": [
                                {"name": "Price objections",
                                 "values": [100, 89, 80, 70]},
                                {"name": "Churn rate",
                                 "values": [100, 123, 145, 159]},
                            ],
                            "source": ("CRM objection log and renewal "
                                       "records, 2023 to 2026")},
                "notes": ("If price were the cause, this line would rise "
                          "with churn. Let the divergence do the work — no "
                          "editorialising needed."),
            },
            {
                "assertion": ("All 10 lost Q4 deals cite the competitor "
                              "guarantee and none cite price"),
                "prose": ("Every lost-deal review from the last quarter "
                          "names the competitor's four-hour response "
                          "guarantee. Price appears in none of them. The "
                          "market has changed what it judges providers on."),
                "bullets": [
                    "Unanimity across 10 reviews is rare — the n is small, the signal is not",
                    "Price at zero mentions rules it out as the loss driver",
                    "The competitor set the frame; the market has adopted it",
                ],
                "source": "Lost-deal reviews, Oct to Dec 2025, n=10",
                "exhibit": {"exhibit_id": "EX-LOSTDEAL",
                            "kind": "bar",
                            "unit": "Mentions across 10 lost-deal reviews",
                            "categories": ["Competitor response guarantee",
                                           "Service scope",
                                           "Relationship change",
                                           "Price"],
                            "values": [10, 3, 2, 0],
                            "highlight": 0,
                            "source": "Lost-deal reviews, Oct to Dec 2025, n=10"},
                "notes": ("This is the exhibit the sceptics remember. Ten "
                          "out of ten is unusual — say so, and say the n is "
                          "small but the signal is unanimous."),
            },
        ],
        "transition": ("If price is not the lever, the lost-deal records say "
                       "what is — and it is buildable."),
    },
    {
        "tag": "REMEDY",
        "assertion": ("A four-hour response guarantee wins the frame the "
                      "market now uses"),
        "why_it_matters": ("The competitor changed the question renewals are "
                           "judged on. Matching the guarantee answers that "
                           "question directly — and Meridian starts closer "
                           "to it than the sales narrative admits."),
        "supports": [
            {
                "assertion": ("The competitor guarantee reset what buyers ask "
                              "at renewal"),
                "prose": ("Renewal conversations now open with response "
                          "commitments, not price. Account managers report "
                          "the guarantee is raised by the customer in most "
                          "competitive renewals. The frame has moved and the "
                          "proposition has not moved with it."),
                "bullets": [
                    "Renewal conversations now open with response commitments",
                    "Customers raise the guarantee unprompted",
                    "The proposition has not moved with the frame",
                ],
                "source": "Account team renewal debriefs, Q1 to Q2 2026",
                "exhibit": None,
                "notes": ("Anchor this in the account managers' own words — "
                          "the evidence is qualitative and honest about it."),
            },
            {
                "assertion": ("Meridian already meets four hours on 78% of "
                              "callouts"),
                "prose": ("Dispatch records show 78% of 2025 callouts were "
                          "answered inside four hours without any guarantee "
                          "in place. The gap to close is the remaining 22% — "
                          "evenings, weekends and multi-site jobs. That gap "
                          "is known and priceable."),
                "bullets": [
                    "This is baseline performance — achieved with no guarantee in place",
                    "The gap is evenings, weekends and multi-site jobs",
                    "The guarantee prices the last 22%, not the whole operation",
                ],
                "source": "Dispatch system records, Jan to Dec 2025, n=8,400",
                "exhibit": {"exhibit_id": "EX-CALLOUTS",
                            "kind": "bar",
                            "unit": "% of 2025 callouts",
                            "categories": ["Met within four hours today",
                                           "Requires new cover"],
                            "values": [78, 22],
                            "highlight": 0,
                            "source": "Dispatch system records, 2025, n=8,400"},
                "notes": ("The surprise is how close we already are. The "
                          "guarantee prices the last 22%, not the whole "
                          "operation."),
            },
            {
                "assertion": ("The guarantee costs £350k a year to stand up"),
                "prose": ("Three cost lines close the 22% gap: extended "
                          "engineering cover, a triage desk, and a service-"
                          "credit reserve priced at a 2% breach rate. The "
                          "total is £350k a year, fixed."),
                "bullets": [
                    "All three lines are fixed cost — no per-callout scaling",
                    "The reserve is the honest line: it prices the 2% breach rate",
                    "£350k is 7% of the £5.1M revenue under protection",
                ],
                "source": "Operations costing model v2, Jul 2026",
                "exhibit": {"exhibit_id": "EX-COSTS",
                            "kind": "waterfall",
                            "unit": "£k a year",
                            "categories": ["Engineering cover", "Triage desk",
                                           "Credit reserve", "Total"],
                            "values": [190, 95, 65, 350],
                            "totals": [3],
                            "source": "Operations costing model v2, Jul 2026"},
                "notes": ("The reserve line is the honest one — a guarantee "
                          "without a priced breach rate is a slogan."),
            },
        ],
        "transition": ("A guarantee Meridian can meet at a known cost must "
                       "still pay. It does — on pessimistic assumptions."),
    },
    {
        "tag": "ECONOMICS",
        "assertion": ("The programme pays back within four quarters on the "
                      "central case"),
        "why_it_matters": ("The case is priced in ranges, not a point value. "
                           "Even the pessimistic case clears the outlay "
                           "inside five quarters, and the staged spend caps "
                           "the downside."),
        "supports": [
            {
                "assertion": ("Recovered renewals return £0.9M to £1.6M a "
                              "year against the £350k outlay"),
                "prose": ("The scenario table runs churn recovery at three "
                          "levels. The central case recovers churn to 27% "
                          "and returns £1.2M a year — payback inside four "
                          "quarters. The pessimistic case still pays back "
                          "in five."),
                "bullets": [
                    "Central case: churn to 27%, £1.2M a year recovered",
                    "Pessimistic case still pays back in five quarters",
                    "All three cases clear the £350k outlay",
                ],
                "source": "Renewal economics model, scenario run Jul 2026",
                "exhibit": {"exhibit_id": "EX-SCENARIOS",
                            "kind": "table",
                            "headers": ["Scenario", "Churn recovered to",
                                        "Revenue recovered", "Payback"],
                            "rows": [
                                ["Pessimistic", "31%", "£0.9M a year", "5 quarters"],
                                ["Central", "27%", "£1.2M a year", "4 quarters"],
                                ["Optimistic", "25%", "£1.6M a year", "3 quarters"],
                            ],
                            "source": "Renewal economics model, Jul 2026"},
                "notes": ("Never present the central case alone. The range "
                          "is the credibility."),
            },
            {
                "assertion": ("The 90-day gate caps the downside at £120k"),
                "prose": ("Spend is staged: £120k reaches the field before "
                          "the 90-day review. If lost-deal citations of the "
                          "competitor guarantee have not fallen below five a "
                          "quarter by then, the programme stops before the "
                          "remaining spend commits."),
                "bullets": [
                    "Only £120k committed before the 90-day review",
                    "Stop trigger: citations still above five a quarter",
                    "The remaining spend never commits on a failed signal",
                ],
                "source": "Programme spend schedule, draft Jul 2026",
                "exhibit": None,
                "notes": ("The gate is what makes this approvable — the "
                          "decision is reversible at a known cost."),
            },
        ],
        "transition": None,
    },
]

CONCLUSION = {
    "restate": ("The renewal problem is a framing problem, and the guarantee "
                "answers it — churn back to 25% and £1.2M of revenue "
                "recovered within four quarters."),
    "key_assumption": ("This rests on the market's response-time frame "
                       "persisting. If lost-deal citations shift away from "
                       "it, we review at the 90-day gate."),
    "next_action": ("The operations director stands up the build team on "
                    "1 Oct 2026."),
}

ASK = {
    "decision": ("Approve the response-time guarantee programme at £350k a "
                 "year"),
    "gate": ("Stop at 90 days if lost-deal citations of the competitor "
             "guarantee have not fallen below five a quarter"),
    "owner": "Operations director",
    "date": "1 Oct 2026",
}

NEXT_STEPS = [
    ("Confirm the service-credit terms with legal", "Commercial lead", "12 Sep 2026"),
    ("Brief the account teams on the renewal script", "Sales director", "19 Sep 2026"),
]

# --- EXHIBIT REUSE (Stage 4) — one row per exhibit this deck uses. Every
# support exhibit's exhibit_id must appear here. decision: "REUSE" (as-is
# from a prior deck) | "ADAPT" ('changes' required) | "NEW". data_vintage =
# date of the UNDERLYING data (YYYY-MM), not the deck date. Staleness guard:
# REUSE/ADAPT rows older than MAX_VINTAGE_MONTHS fail the build unless
# vintage_accepted is True with a stated vintage_reason. NEW rows are exempt.
# ENGAGEMENT ASSET: never commit a filled register to the public repository.
MAX_VINTAGE_MONTHS = 6

EXHIBIT_REGISTER = [
    {"exhibit_id": "EX-OBJECTIONS",
     "message": ("Price objections fell while churn rose — the two series "
                 "diverge"),
     "chart_form": "line", "source_deck": "Q2 renewal review deck",
     "data_vintage": "2026-07", "decision": "ADAPT",
     "changes": ("Add the churn series alongside objections and index both "
                 "to 2023 = 100 so the divergence reads on one axis"),
     "vintage_accepted": False, "vintage_reason": ""},
    {"exhibit_id": "EX-LOSTDEAL",
     "message": ("All 10 lost Q4 deals cite the competitor guarantee; none "
                 "cite price"),
     "chart_form": "bar", "source_deck": "",
     "data_vintage": "2025-12", "decision": "NEW", "changes": "",
     "vintage_accepted": False, "vintage_reason": ""},
    {"exhibit_id": "EX-CALLOUTS",
     "message": "Meridian already meets four hours on 78% of callouts",
     "chart_form": "bar", "source_deck": "",
     "data_vintage": "2026-02", "decision": "NEW", "changes": "",
     "vintage_accepted": False, "vintage_reason": ""},
    {"exhibit_id": "EX-COSTS",
     "message": "The guarantee costs £350k a year, in three fixed lines",
     "chart_form": "waterfall", "source_deck": "Operations costing pack v2",
     "data_vintage": "2026-01", "decision": "ADAPT",
     "changes": ("Rebuild the costing table as a build-up waterfall so the "
                 "three lines read as one £350k total"),
     "vintage_accepted": True,
     "vintage_reason": ("Costing re-validated Jul 2026 against supplier "
                        "quotes; underlying rates unchanged")},
    {"exhibit_id": "EX-SCENARIOS",
     "message": ("Recovered renewals return £0.9M–£1.6M a year against the "
                 "£350k outlay"),
     "chart_form": "table", "source_deck": "Renewal economics model pack",
     "data_vintage": "2026-07", "decision": "REUSE", "changes": "",
     "vintage_accepted": False, "vintage_reason": ""},
]

# ============================================================================
# BUILD — DO NOT EDIT BELOW THIS LINE
# ============================================================================
import datetime
import re
import sys

OUTPUT_FILE = "standalone-deck.pptx"

# ---------------------------------------------------------------- self-check
BANNED_LABELS = {"background", "analysis", "overview", "introduction",
                 "context", "summary", "next steps", "recommendation",
                 "findings", "update", "situation", "options", "approach",
                 "details", "conclusion"}

def _assertion_defect(text):
    t = (text or "").strip()
    words = t.split()
    if len(words) < 5:
        return "under 5 words — reads as a label, not a claim"
    if len(words) > 15:
        return "over 15 words — will not fit the title zone; sharpen the claim"
    if t.rstrip(".").lower() in BANNED_LABELS:
        return "banned label heading"
    if t.endswith(":"):
        return "ends in a colon — a label, not a claim"
    tail = [w for w in words[1:] if w[:1].isalpha()]
    caps = sum(1 for w in tail if w[:1].isupper())
    if tail and caps > 0.6 * len(tail):
        return "Title Case — use sentence case"
    return None

def _collect_strings(obj, out):
    if isinstance(obj, str):
        out.append(obj)
    elif isinstance(obj, dict):
        for v in obj.values():
            _collect_strings(v, out)
    elif isinstance(obj, (list, tuple)):
        for v in obj:
            _collect_strings(v, out)

def _valid_exhibit(ex):
    if ex is None:
        return None
    if not isinstance(ex, dict):
        return "exhibit is not a dict"
    kind = ex.get("kind")
    if kind in ("bar", "column", "line"):
        cats = ex.get("categories")
        if not cats:
            return "chart has no categories"
        series = ex.get("series")
        if series is not None:
            # multi-series: [{"name": ..., "values": [...]}, ...], max two
            if not (isinstance(series, list) and 1 <= len(series) <= 2):
                return "series must hold one or two {name, values} entries"
            for sr in series:
                if not (isinstance(sr, dict) and sr.get("name")
                        and sr.get("values")):
                    return "each series needs a name and values"
                if len(sr["values"]) != len(cats):
                    return ("series '%s' length does not match categories"
                            % sr["name"])
            if len(series) > 1 and ex.get("highlight") is not None:
                return "highlight applies to single-series charts only"
        else:
            vals = ex.get("values")
            if not vals or len(cats) != len(vals):
                return "chart categories/values missing or mismatched"
            hl = ex.get("highlight")
            if hl is not None and not (0 <= hl < len(vals)):
                return "highlight index out of range"
    elif kind == "waterfall":
        cats, vals = ex.get("categories"), ex.get("values")
        if not cats or not vals or len(cats) != len(vals):
            return "waterfall categories/values missing or mismatched"
        totals = ex.get("totals", [])
        if not all(isinstance(t, int) and 0 <= t < len(vals)
                   for t in totals):
            return "waterfall totals indices out of range"
        if not all(isinstance(v, (int, float)) for v in vals):
            return "waterfall values must be numbers (signed deltas)"
    elif kind == "table":
        if not ex.get("headers") or not ex.get("rows"):
            return "table needs headers and rows"
    else:
        return "exhibit kind must be bar / column / line / waterfall / table"
    if not ex.get("source"):
        return "exhibit has no source"
    return None

EXHIBIT_DECISIONS = ("REUSE", "ADAPT", "NEW")

def _theme_defects():
    defects = []
    for k, v in THEME.items():
        if k.startswith("font_"):
            if not str(v).strip():
                defects.append("THEME['%s'] is empty — name a real font" % k)
        elif not re.fullmatch(r"[0-9A-Fa-f]{6}", str(v)):
            defects.append("THEME['%s'] = %r is not a 6-char hex "
                           "(no '#', no alpha)" % (k, v))
    return defects

def _vintage_age_months(vintage):
    """Age of a YYYY-MM vintage in months, or None if malformed."""
    m = re.fullmatch(r"(\d{4})-(\d{2})", str(vintage).strip())
    if not m:
        return None
    y, mo = int(m.group(1)), int(m.group(2))
    if not 1 <= mo <= 12:
        return None
    today = datetime.date.today()
    return (today.year - y) * 12 + (today.month - mo)

def _exhibit_register_defects():
    """Structural defects and the staleness guard, separately named."""
    struct, stale = [], []
    seen = set()
    for r in EXHIBIT_REGISTER:
        rid = str(r.get("exhibit_id", "")).strip() or "?"
        if rid == "?":
            struct.append("a register row has no exhibit_id")
        elif rid in seen:
            struct.append("duplicate exhibit_id %s" % rid)
        seen.add(rid)
        if not str(r.get("message", "")).strip():
            struct.append("%s has no message — name the claim it proves"
                          % rid)
        dec = r.get("decision")
        if dec not in EXHIBIT_DECISIONS:
            struct.append("%s decision must be REUSE / ADAPT / NEW" % rid)
        if dec == "ADAPT" and not str(r.get("changes", "")).strip():
            struct.append("%s is ADAPT with no 'changes' — say what "
                          "changes" % rid)
        if r.get("vintage_accepted") and not str(
                r.get("vintage_reason", "")).strip():
            struct.append("%s is vintage_accepted with no stated reason"
                          % rid)
        age = _vintage_age_months(r.get("data_vintage", ""))
        if dec in ("REUSE", "ADAPT"):
            if not str(r.get("source_deck", "")).strip():
                struct.append("%s is %s with no source_deck" % (rid, dec))
            if age is None:
                struct.append("%s data_vintage must be YYYY-MM (date of "
                              "the underlying data)" % rid)
            elif age > MAX_VINTAGE_MONTHS and not (
                    r.get("vintage_accepted")
                    and str(r.get("vintage_reason", "")).strip()):
                stale.append("%s: data vintage %s is %d months old (max "
                             "%d) — rebuild as NEW or set vintage_accepted "
                             "with a reason"
                             % (rid, r["data_vintage"], age,
                                MAX_VINTAGE_MONTHS))
    return struct, stale

def run_self_check():
    checks, warnings = [], []
    td = _theme_defects()
    checks.append(("THEME: colours are 6-char hex and fonts are named"
                   + ("" if not td else " — " + "; ".join(td)), not td))
    gt_words = len(GOVERNING_THOUGHT.split())
    checks.append(("governing thought is 25 words or fewer (%d)" % gt_words,
                   0 < gt_words <= 25))
    checks.append(("SCQA complete: situation, complication, question, answer",
                   all([SCQA.get("situation"), SCQA.get("complication"),
                        SCQA.get("question"), GOVERNING_THOUGHT])))
    low_sit = " " + SCQA["situation"].lower()
    checks.append(("situation carries no 'but' / 'however'",
                   " but " not in low_sit and "however" not in low_sit))
    checks.append(("exactly one question, ending in a question mark",
                   SCQA["question"].count("?") == 1
                   and SCQA["question"].strip().endswith("?")))
    checks.append(("key line holds 2–4 arguments (%d)" % len(KEY_LINE),
                   2 <= len(KEY_LINE) <= 4))
    checks.append(("ordering principle named: time / structure / degree",
                   KEY_LINE_ORDER in ("time", "structure", "degree")))

    label_defects = []
    for kl in KEY_LINE:
        d = _assertion_defect(kl.get("assertion", ""))
        if d:
            label_defects.append("key line '%s…': %s"
                                 % (kl.get("assertion", "")[:40], d))
        for sp in kl.get("supports", []):
            d = _assertion_defect(sp.get("assertion", ""))
            if d:
                label_defects.append("support '%s…': %s"
                                     % (sp.get("assertion", "")[:40], d))
    checks.append(("every slide title is a full-sentence claim, not a label"
                   + ("" if not label_defects else
                      " — " + "; ".join(label_defects)),
                   not label_defects))

    checks.append(("every key-line argument carries 2–5 supports",
                   all(2 <= len(kl.get("supports", [])) <= 5
                       for kl in KEY_LINE)))
    kl_assertions = [kl["assertion"].strip().lower() for kl in KEY_LINE]
    checks.append(("no two key-line arguments make the same claim (MECE "
                   "overlap)", len(set(kl_assertions)) == len(kl_assertions)))

    # standalone discipline: the slide itself must prove the title — and fit
    prose_defects = []
    for kl in KEY_LINE:
        for sp in kl["supports"]:
            n = len((sp.get("prose") or "").split())
            if n < 15:
                prose_defects.append("'%s…': on-slide prose under 15 words — "
                                     "a cold reader cannot follow"
                                     % sp["assertion"][:35])
            elif n > 55:
                prose_defects.append("'%s…': on-slide prose is %d words "
                                     "(max 55) — it will collide with the "
                                     "exhibit; cut or split the slide"
                                     % (sp["assertion"][:35], n))
    checks.append(("standalone bodies: every support carries on-slide prose "
                   "of 15–55 words"
                   + ("" if not prose_defects else
                      " — " + "; ".join(prose_defects)),
                   not prose_defects))

    # geometry guards: everything must fit its fixed zone
    geo_defects = []
    for kl in KEY_LINE:
        wm = len((kl.get("why_it_matters") or "").split())
        if wm > 40:
            geo_defects.append("'%s…': why_it_matters is %d words (max 40) "
                               "— it overflows the divider"
                               % (kl["assertion"][:35], wm))
        for sp in kl["supports"]:
            nb = len(sp.get("bullets") or [])
            if not 1 <= nb <= 3:
                geo_defects.append("'%s…': %d bullets — the slide shows "
                                   "exactly one to three; cut to the ones "
                                   "that matter" % (sp["assertion"][:35], nb))
    checks.append(("geometry: why_it_matters 40 words max, 1–3 bullets per "
                   "support"
                   + ("" if not geo_defects else
                      " — " + "; ".join(geo_defects)), not geo_defects))

    # a support that argues with numbers should show them, not say them
    for kl in KEY_LINE:
        for sp in kl["supports"]:
            if sp.get("exhibit") is None and re.search(
                    r"\d+%|£[\d,.]+|\d{2,}",
                    " ".join(sp.get("bullets", []) or []) + " "
                    + (sp.get("prose") or "")):
                warnings.append("'%s…' quotes numbers but carries no "
                                "exhibit — consider a chart"
                                % sp["assertion"][:35])
    checks.append(("every evidence slide cites its source on the slide",
                   all(sp.get("source") for kl in KEY_LINE
                       for sp in kl["supports"])))

    ex_defects = []
    for kl in KEY_LINE:
        for sp in kl["supports"]:
            d = _valid_exhibit(sp.get("exhibit"))
            if d:
                ex_defects.append("'%s…': %s" % (sp["assertion"][:35], d))
    checks.append(("every exhibit is well-formed with a source"
                   + ("" if not ex_defects else " — " + "; ".join(ex_defects)),
                   not ex_defects))

    # --- exhibit reuse register (Stage 4)
    reg_struct, reg_stale = _exhibit_register_defects()
    checks.append(("exhibit register rows are well-formed: message, valid "
                   "decision, ADAPT changes, REUSE/ADAPT source + vintage"
                   + ("" if not reg_struct else
                      " — " + "; ".join(reg_struct)), not reg_struct))
    checks.append(("staleness guard: no REUSE/ADAPT older than %d months "
                   "without an accepted vintage" % MAX_VINTAGE_MONTHS
                   + ("" if not reg_stale else
                      " — " + "; ".join(reg_stale)), not reg_stale))
    reg_ids = set(str(r.get("exhibit_id", "")).strip()
                  for r in EXHIBIT_REGISTER)
    link_defects = []
    for kl in KEY_LINE:
        for sp in kl["supports"]:
            ex = sp.get("exhibit")
            if ex is None:
                continue
            eid = str(ex.get("exhibit_id", "")).strip()
            if not eid:
                link_defects.append("'%s…' exhibit has no exhibit_id"
                                    % sp["assertion"][:35])
            elif eid not in reg_ids:
                link_defects.append("'%s…' exhibit_id %s has no register "
                                    "row" % (sp["assertion"][:35], eid))
    checks.append(("every support exhibit carries an exhibit_id with a "
                   "register row (REUSE / ADAPT / NEW decided)"
                   + ("" if not link_defects else
                      " — " + "; ".join(link_defects)), not link_defects))

    checks.append(("conclusion has Minto's three parts: restate, key "
                   "assumption, next action",
                   all([CONCLUSION.get("restate"),
                        CONCLUSION.get("key_assumption"),
                        CONCLUSION.get("next_action")])))
    checks.append(("single ask with decision, gate, owner and date",
                   all([ASK.get("decision"), ASK.get("gate"),
                        ASK.get("owner"), ASK.get("date")])))

    strings = []
    for blob in (META, SCQA, GOVERNING_THOUGHT, KEY_LINE, CONCLUSION, ASK,
                 NEXT_STEPS):
        _collect_strings(blob, strings)
    placeholders = sorted(set(m for s in strings
                              for m in re.findall(r"\[[^\]]+\]", s)))
    checks.append(("no bracketed placeholders remain"
                   + ("" if not placeholders else
                      " — replace: " + ", ".join(placeholders)),
                   not placeholders))
    return checks, warnings

def print_report(checks, warnings, extra=(), wrote=False):
    print("=" * 64)
    print("SELF-CHECK — %s" % OUTPUT_FILE)
    for name, ok in list(checks) + list(extra):
        print("  [%s] %s" % ("PASS" if ok else "FAIL", name))
    for w in warnings:
        print("  [WARN] %s" % w)
    print("=" * 64)
    if not all(ok for _, ok in list(checks) + list(extra)):
        raise SystemExit("SELF-CHECK FAILED — fix the DATA section, not the "
                         "checks.")
    if wrote:
        print("written", OUTPUT_FILE)

_checks, _warnings = run_self_check()
if not all(ok for _, ok in _checks):
    print_report(_checks, _warnings)
    sys.exit(1)

# ---------------------------------------------------------------- pptx build
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION,
                             XL_LEGEND_POSITION)
from pptx.oxml.ns import qn

def _rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _blend(base_hex, mix_hex, f):
    """f = share of mix_hex blended into base_hex."""
    b = [int(base_hex[i:i + 2], 16) for i in (0, 2, 4)]
    m = [int(mix_hex[i:i + 2], 16) for i in (0, 2, 4)]
    return RGBColor(*(round(b[i] * (1 - f) + m[i] * f) for i in range(3)))

SLATE = _rgb(THEME["dominant"])
BLUE = _rgb(THEME["accent"])
GOLD = _rgb(THEME["highlight"])
TINT = _rgb(THEME["card_fill"])
INK = _rgb(THEME["ink"])
BODY = _rgb(THEME["ink_soft"])
MUTE = _rgb(THEME["muted"])
GREYC = _rgb(THEME["chart_compare"])
WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # functional — not brand
# derived tints of the dominant colour, for text on dark slides
DOM_LT = _blend(THEME["dominant"], "FFFFFF", 0.75)   # light body on dark
DOM_MID = _blend(THEME["dominant"], "FFFFFF", 0.50)  # metadata on dark
DOM_NUM = _blend(THEME["dominant"], "FFFFFF", 0.22)  # watermark numerals
HEAD_FONT, BODY_FONT = THEME["font_heading"], THEME["font_body"]

W, H, M = 13.333, 7.5, 0.62
CW = W - 2 * M

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = fill   # invisible border: line colour = fill colour
    shp.line.width = Pt(0)
    if rounded:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp

def flood(s, fill):
    return rect(s, -0.02, -0.02, W + 0.04, H + 0.04, fill)

def txt(s, x, y, w, h, runs, size=13, color=BODY, bold=False, italic=False,
        align=PP_ALIGN.LEFT, spacing=None, font=BODY_FONT, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for text, opt in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opt.get("align", align)
        if spacing:
            p.line_spacing = Pt(spacing)
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = opt.get("font", font)
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", bold)
        f.italic = opt.get("italic", italic)
        f.color.rgb = opt.get("color", color)
    return tb

def bullets(s, items, x, y, w, size=13, color=BODY, gap=8):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "•  " + it
        f = r.font
        f.name = BODY_FONT
        f.size = Pt(size)
        f.color.rgb = color
    return tb

def eyebrow(s, t, y=0.32):
    txt(s, M, y, 10, 0.3, t.upper(), size=11, color=GOLD, bold=True)

def title(s, t, size=22, color=BLUE):
    txt(s, M, 0.66, CW, 1.0, t, size=size, color=color, font=HEAD_FONT,
        spacing=size + 5, anchor=MSO_ANCHOR.TOP)

def foot(s, t):
    txt(s, M, H - 0.55, CW, 0.3, t, size=9, color=MUTE)

def pageno(s):
    txt(s, W - M - 0.6, H - 0.55, 0.6, 0.3,
        "%02d" % len(prs.slides._sldIdLst), size=9, color=MUTE,
        align=PP_ALIGN.RIGHT)

def notes(s, t):
    s.notes_slide.notes_text_frame.text = t

def _strip_axes(ch):
    va = ch.value_axis
    va.visible = False
    va.has_major_gridlines = False
    ca = ch.category_axis
    ca.format.line.fill.background()
    ca.tick_labels.font.size = Pt(11)
    ca.tick_labels.font.name = BODY_FONT

def _hide_series_labels(ser):
    """Series-level <c:dLbls><c:delete/> — overrides plot-level labels."""
    el = ser._element
    dLbls = el.makeelement(qn("c:dLbls"), {})
    dLbls.append(el.makeelement(qn("c:delete"), {"val": "1"}))
    ref = el.find(qn("c:cat"))
    if ref is None:
        ref = el.find(qn("c:val"))
    if ref is not None:
        ref.addprevious(dLbls)
    else:
        el.append(dLbls)

def _add_waterfall(s, ex, x, y, w, h):
    """Stacked column with an invisible base series. values are signed
    deltas; indices in ex["totals"] are absolute totals drawn from zero.
    Deltas colour accent (up) / comparison grey (down); totals dominant."""
    vals, totals = ex["values"], set(ex.get("totals", []))
    bases, bars = [], []
    running = 0.0
    for i, v in enumerate(vals):
        if i in totals:
            bases.append(0)
            bars.append(v)
            running = v
        else:
            bases.append(running if v >= 0 else running + v)
            bars.append(abs(v))
            running += v
    cd = CategoryChartData()
    cd.categories = [str(c) for c in ex["categories"]]
    cd.add_series("base", tuple(bases))
    cd.add_series(ex.get("unit", "value"), tuple(bars))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(x),
                            Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    plot = ch.plots[0]
    plot.gap_width = 50
    base_ser, bar_ser = plot.series[0], plot.series[1]
    base_ser.format.fill.background()
    base_ser.format.line.fill.background()
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(12)
    dl.font.name = BODY_FONT
    dl.font.bold = True
    dl.font.color.rgb = WHITE
    dl.position = XL_LABEL_POSITION.INSIDE_END
    _hide_series_labels(base_ser)
    for i in range(len(vals)):
        pt = bar_ser.points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = (SLATE if i in totals
                                         else BLUE if vals[i] >= 0
                                         else GREYC)
    _strip_axes(ch)
    return gf

def add_chart(s, ex, x, y, w, h):
    if ex["kind"] == "waterfall":
        return _add_waterfall(s, ex, x, y, w, h)
    kinds = {"bar": XL_CHART_TYPE.BAR_CLUSTERED,
             "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
             "line": XL_CHART_TYPE.LINE}
    cd = CategoryChartData()
    cd.categories = [str(c) for c in ex["categories"]]
    series_spec = ex.get("series")
    if series_spec:
        for sr in series_spec:
            cd.add_series(sr["name"], tuple(sr["values"]))
    else:
        cd.add_series(ex.get("unit", "value"), tuple(ex["values"]))
    gf = s.shapes.add_chart(kinds[ex["kind"]], Inches(x), Inches(y),
                            Inches(w), Inches(h), cd)
    ch = gf.chart
    multi = bool(series_spec) and len(series_spec) > 1
    ch.has_legend = multi
    if multi:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(11)
        ch.legend.font.name = BODY_FONT
        ch.legend.font.color.rgb = BODY
    ch.has_title = False
    plot = ch.plots[0]
    if ex["kind"] != "line":
        plot.gap_width = 60
    if multi:
        # legend carries the meaning; labels would clutter two series
        if ex["kind"] == "line":
            colors = (BLUE, GOLD)
            for i, ser in enumerate(plot.series):
                ser.format.line.color.rgb = colors[i % 2]
                ser.format.line.width = Pt(2.5)
                ser.smooth = False
        else:
            colors = (SLATE, GREYC)
            for i, ser in enumerate(plot.series):
                ser.format.fill.solid()
                ser.format.fill.fore_color.rgb = colors[i % 2]
    else:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(12)
        dl.font.name = BODY_FONT
        dl.font.bold = True
        dl.font.color.rgb = INK
        if ex["kind"] != "line":
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
            ser = plot.series[0]
            hl = ex.get("highlight")
            for i in range(len(ex["values"])):
                pt = ser.points[i]
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = (SLATE if hl is not None
                                                 and i == hl else GREYC)
        else:
            ser = plot.series[0]
            ser.format.line.color.rgb = BLUE
            ser.format.line.width = Pt(2.5)
    _strip_axes(ch)
    return gf

def add_table(s, ex, x, y, w):
    headers, rows = ex["headers"], ex["rows"]
    nrows, ncols = 1 + len(rows), len(headers)
    row_h = 0.48
    shp = s.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w),
                             Inches(row_h * nrows))
    tbl = shp.table
    for j, hd in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLATE
        cell.text = str(hd)
        pr = cell.text_frame.paragraphs[0].runs[0].font
        pr.size = Pt(12)
        pr.bold = True
        pr.name = HEAD_FONT
        pr.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell = tbl.cell(1 + i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else TINT
            cell.text = str(v)
            pr = cell.text_frame.paragraphs[0].runs[0].font
            pr.size = Pt(12)
            pr.name = BODY_FONT
            pr.color.rgb = INK
            pr.bold = (i == len(rows) - 1
                       and str(row[0]).lower().startswith("total"))
    return shp

# ---------------- 1. COVER ----------------
s = slide()
flood(s, SLATE)
rect(s, 0, 0, W, 0.16, GOLD)
txt(s, M, 1.5, 10.5, 0.4, META["subtitle"].upper(), size=13, color=GOLD,
    bold=True)
txt(s, M, 2.05, 11.6, 2.2, META["title"], size=44, color=WHITE,
    font=HEAD_FONT, spacing=50)
txt(s, M, 4.35, 11.2, 0.9, "A standalone read — the full case, evidence and "
    "the decision requested, in %d pages." % (
        3 + sum(1 + len(kl["supports"]) for kl in KEY_LINE) + 2),
    size=15, color=DOM_LT, spacing=21)
txt(s, M, 6.5, 10, 0.4, "%s   ·   %s   ·   for: %s" % (
    META["author_line"], META["date"], META["audience"]), size=11.5,
    color=DOM_MID)
notes(s, "Standalone deck — reads unaided. Generated by "
         "generate_standalone_deck.py; edit the DATA section and re-run.")

# ---------------- 2. EXECUTIVE SUMMARY ----------------
s = slide()
eyebrow(s, "Executive summary — the whole case on one page")
title(s, GOVERNING_THOUGHT, size=20, color=SLATE)
y = 2.0
for lab, text in (("SITUATION", SCQA["situation"]),
                  ("COMPLICATION", SCQA["complication"]),
                  ("QUESTION", SCQA["question"])):
    rect(s, M, y, CW, 1.02, TINT, rounded=True)
    txt(s, M + 0.4, y + 0.12, 2.4, 0.3, lab, size=10.5, color=BLUE, bold=True)
    txt(s, M + 0.4, y + 0.42, CW - 0.9, 0.55, text, size=12, color=BODY,
        spacing=15)
    y += 1.16
rect(s, M, y, CW, 1.35, SLATE, rounded=True)
txt(s, M + 0.4, y + 0.12, 3.5, 0.3, "THE ANSWER, IN %d ARGUMENTS"
    % len(KEY_LINE), size=10.5, color=GOLD, bold=True)
txt(s, M + 0.4, y + 0.42,
    CW - 0.9, 0.85, "   ·   ".join("%d) %s" % (i + 1, kl["assertion"])
                                   for i, kl in enumerate(KEY_LINE)),
    size=12, color=WHITE, spacing=16)
pageno(s)
notes(s, "The 30-second read: SCQA plus the key line. A reader who stops "
         "here has the whole argument.")

# ---------------- 3. AGENDA ----------------
s = slide()
eyebrow(s, "How the case is laid out")
title(s, "The argument runs in %s order" % KEY_LINE_ORDER)
y = 2.1
for i, kl in enumerate(KEY_LINE):
    rect(s, M, y, CW, 1.05, TINT, rounded=True)
    txt(s, M + 0.4, y + 0.26, 0.8, 0.55, str(i + 1), size=26, color=GOLD,
        bold=True, font=HEAD_FONT)
    txt(s, M + 1.3, y + 0.14, 3.0, 0.3, kl["tag"].upper(), size=10.5,
        color=BLUE, bold=True)
    txt(s, M + 1.3, y + 0.44, CW - 2.0, 0.55, kl["assertion"], size=14.5,
        color=INK, bold=True, font=HEAD_FONT, spacing=18)
    y += 1.22
txt(s, M, y + 0.05, CW, 0.5, "Then the conclusion, the key assumption, and "
    "the single decision requested.", size=12, color=MUTE, italic=True)
pageno(s)
notes(s, "Agenda doubles as the title-read test: the three claims alone "
         "carry the argument.")

# ---------------- SECTIONS ----------------
for i, kl in enumerate(KEY_LINE):
    s = slide()
    flood(s, SLATE)
    txt(s, M, 2.0, 2.4, 1.6, "%02d" % (i + 1), size=80,
        color=DOM_NUM, bold=True, font=HEAD_FONT)
    txt(s, M + 2.45, 2.18, 9.6, 0.4, kl["tag"].upper(), size=12, color=GOLD,
        bold=True)
    txt(s, M + 2.45, 2.6, 9.6, 1.4, kl["assertion"], size=30, color=WHITE,
        font=HEAD_FONT, spacing=37)
    txt(s, M + 2.45, 4.05, 9.4, 0.9, kl["why_it_matters"], size=13.5,
        color=DOM_LT, spacing=19)
    # section position strip
    strip = "   ·   ".join(("■ " if j == i else "□ ") + k2["tag"].title()
                           for j, k2 in enumerate(KEY_LINE))
    txt(s, M + 2.45, 6.4, 9.4, 0.35, strip, size=11,
        color=DOM_MID)
    notes(s, kl.get("transition") or "Section divider.")

    for sp in kl["supports"]:
        s = slide()
        eyebrow(s, "%s  ·  %02d" % (kl["tag"], i + 1))
        title(s, sp["assertion"], size=21)
        # standalone: prose proves the title on the slide itself
        txt(s, M, 1.72, CW, 1.0, sp["prose"], size=13, color=BODY,
            spacing=18)
        ex = sp.get("exhibit")
        if ex and ex["kind"] in ("bar", "column", "line", "waterfall"):
            txt(s, M, 2.62, 7.6, 0.3, ex.get("unit", ""), size=11, color=MUTE)
            add_chart(s, ex, M, 2.94, 7.6, 3.75)
            rect(s, M + 8.0, 2.94, CW - 8.0, 3.75, TINT, rounded=True)
            txt(s, M + 8.3, 3.2, CW - 8.6, 0.3, "WHAT THIS SHOWS", size=10,
                color=BLUE, bold=True)
            bullets(s, sp["bullets"][:3], M + 8.3, 3.6, CW - 8.65,
                    size=11.5, gap=8)
        elif ex and ex["kind"] == "table":
            add_table(s, ex, M, 2.85, 8.2)
            rect(s, M + 8.6, 2.85, CW - 8.6, 3.6, TINT, rounded=True)
            txt(s, M + 8.9, 3.1, CW - 9.2, 0.3, "WHAT THIS SHOWS", size=10,
                color=BLUE, bold=True)
            bullets(s, sp["bullets"][:3], M + 8.9, 3.5, CW - 9.25,
                    size=11, gap=8)
        else:
            # parallel cards — supports without an exhibit read horizontally
            n = max(1, len(sp["bullets"][:3]))
            card_w = (CW - 0.3 * (n - 1)) / n
            for bi, b in enumerate(sp["bullets"][:3]):
                xx = M + bi * (card_w + 0.3)
                rect(s, xx, 2.9, card_w, 2.7, TINT, rounded=True)
                rect(s, xx, 2.9, card_w, 0.07, GOLD)
                txt(s, xx + 0.32, 3.05, card_w - 0.64, 2.4, b, size=13.5,
                    color=INK, spacing=19, anchor=MSO_ANCHOR.MIDDLE)
        foot(s, "SOURCE: " + sp["source"])
        pageno(s)
        notes(s, sp.get("notes") or "")

# ---------------- CONCLUSION ----------------
s = slide()
eyebrow(s, "Conclusion")
title(s, CONCLUSION["restate"], size=20, color=SLATE)
y = 2.5
for lab, text in (("KEY ASSUMPTION", CONCLUSION["key_assumption"]),
                  ("NEXT ACTION", CONCLUSION["next_action"])):
    rect(s, M, y, CW, 1.25, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.18, 3.4, 0.3, lab, size=10.5, color=BLUE,
        bold=True)
    txt(s, M + 0.42, y + 0.5, CW - 0.9, 0.65, text, size=13.5, color=INK,
        spacing=18)
    y += 1.45
if NEXT_STEPS:
    txt(s, M, y + 0.05, CW, 0.35, "Self-evident next steps", size=12,
        color=SLATE, bold=True, font=HEAD_FONT)
    bullets(s, ["%s — %s, %s" % t for t in NEXT_STEPS], M, y + 0.45, CW,
            size=11.5, gap=5)
pageno(s)
notes(s, "Minto's three-part close: restate, name the assumption, name the "
         "next action.")

# ---------------- ASK ----------------
s = slide()
flood(s, SLATE)
rect(s, 0, H - 0.16, W, 0.16, GOLD)
txt(s, M, 0.9, 10, 0.4, "THE DECISION REQUESTED", size=13, color=GOLD,
    bold=True)
txt(s, M, 1.5, CW, 1.6, ASK["decision"], size=30, color=WHITE,
    font=HEAD_FONT, spacing=38)
y = 3.5
for lab, text in (("Gate", ASK["gate"]), ("Owner", ASK["owner"]),
                  ("Date", ASK["date"])):
    txt(s, M, y, 1.8, 0.4, lab.upper(), size=11, color=GOLD, bold=True)
    txt(s, M + 1.9, y - 0.03, CW - 2.5, 0.6, text, size=14, color=WHITE,
        spacing=19)
    y += 0.75
txt(s, M, 6.55, CW, 0.35, META["footer_tag"], size=9,
    color=DOM_MID)
notes(s, "One decision, one gate, one owner, one date — nothing else on the "
         "page.")

# ---------------------------------------------------------------- post-check
expected = 3 + sum(1 + len(kl["supports"]) for kl in KEY_LINE) + 2
extra = [("slide count matches the pyramid (%d expected)" % expected,
          len(prs.slides._sldIdLst) == expected)]
if all(ok for _, ok in extra):
    prs.save(OUTPUT_FILE)
    print_report(_checks, _warnings, extra, wrote=True)
else:
    print_report(_checks, _warnings, extra, wrote=False)
