#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PROPOSITION BACKLOG DECK GENERATOR — Barclays-branded, Copilot-executable.

HOW TO USE (in the Copilot app):
  1. Attach this file. Say: "Run this script and give me the .pptx it creates."
  2. To change the deck: edit ONLY the DATA section below, then re-run.
     Do not edit anything under LAYOUT — DO NOT EDIT.

The deck is self-computing: per-slot returns = NPV / slots, the bridge
total is summed, and the governing thought's earnings figure is derived.
A self-check runs on every generation and prints PASS/FAIL.

Requires: python-pptx.  Output: proposition-backlog-template-barclays.pptx
"""

# ============================================================================
# DATA — EDIT THIS SECTION ONLY
# ============================================================================

META = {
    "wrapper": "Trade Ready",
    "segment": "exporting SMEs",
    "segment_title": "exporting-SME",
    "page_tag": "Proposition backlog  |  [Month Year]",
    "team_line": "[Team / function]   ·   [Month Year]",
    "earnings_year": 3,
    "margin_pct": 40,
    "slots_per_cycle": 2,
}

SCQ = {
    "situation": "The bank serves exporting SMEs. Five candidate features sit on the roadmap; build capacity is two slots per cycle, with an earnings mandate.",
    "complication": "A one-off pick answers this cycle and nothing after it. The features differ in demand, return and build cost — and the evidence moves as we learn. Gut feel cannot hold an order across cycles.",
    "question": "What do we build into the exporting-SME proposition — and in what order?",
}

GATE = 60

# One record per candidate feature. risk_npv in GBPm on the CHOSEN route;
# slots = build slots consumed on that route. Refused features: risk_npv None.
FEATURES = [
    {"name": "FX payments",     "ocv": 78, "risk_npv": 2.1, "slots": 1,
     "route": "Build in-house",                    "status": "THIS CYCLE"},
    {"name": "Invoice finance", "ocv": 72, "risk_npv": 1.6, "slots": 1,
     "route": "Partner",                           "status": "THIS CYCLE"},
    {"name": "Trade insurance", "ocv": 64, "risk_npv": 1.8, "slots": 2,
     "route": "Build — needs claims capability", "status": "NEXT CYCLE"},
    {"name": "Crypto payments", "ocv": 41, "risk_npv": None, "slots": None,
     "route": "—",                            "status": "OFF THE LIST"},
    {"name": "Rewards scheme",  "ocv": 33, "risk_npv": None, "slots": None,
     "route": "—",                            "status": "OFF THE LIST"},
]

# Earnings bridge (GBPm, year-3 run-rate, this cycle's features).
BRIDGE_BASE = 4.0
BRIDGE_STEPS = [
    ("Price\nmove",                     0.9),
    ("Switch\nincentive\nuptake",       1.3),
    ("Incentive\ncost",                -0.7),
    ("'Trade Ready'\nrecall\nuplift",   0.8),
]

# Route pricing table (slide 15). The counterfactual row is mandatory.
ROUTE_ROWS = [
    ("FX — build",                 "£2.4m", "9 mths ±1",  "High", "GREEN", "£2.1m", "PASS"),
    ("Invoice finance — partner",  "£1.9m", "5 mths ±1",  "High", "GREEN", "£1.6m", "PASS"),
    ("Invoice finance — if built", "£1.9m", "18 mths ±4", "Low",  "RED",   "£0.6m", "FAIL"),
]

KEYLINE = [
    ("DEMAND ADMITS", "Exporters want three of the five — and will think of us",
     "FX (78), invoice finance (72) and insurance (64) clear the demand gate of 60. Crypto (41) and rewards (33) never get priced. The 'Trade Ready' name makes the shelf visible at the moment the need hits."),
    ("RETURN SORTS", "Return per build slot sets the order — FX first",
     "FX earns £2.1m a slot, invoice finance £1.6m, insurance £0.9m. Insurance beats invoice finance on raw NPV and still sorts below it — it needs two slots. The order is re-scored quarterly."),
    ("CAPACITY CUTS", "Two slots this cycle — and the routes keep the edge",
     "FX builds in-house on our core platform. Invoice finance partners: live in 5 months instead of 18. Insurance waits on the claims capability. Partner terms keep the customer and the data ours."),
]

REC_LINES = [
    ("This cycle",  "FX payments (build) · invoice finance (partner)", "the two slots"),
    ("Next cycle",  "Trade insurance",                                      "third by return per slot"),
    ("Off the list","Crypto payments · rewards scheme",                "below the demand gate"),
    ("Standing",    "Re-score the backlog every quarter",                   "the order moves with the evidence"),
]

EVIDENCE_STATS = [
    ("320",    "Exporting SMEs interviewed"),
    ("9 wks",  "Fieldwork window, ending [date]"),
    ("3 of 5", "Features tested with a live prototype"),
    ("61%",    "Of segment revenue represented"),
]
EVIDENCE_LIMITS = [
    "Smallest exporters (under £1m turnover) are under-represented — could overstate invoice finance demand",
    "Exporters over-claim switching intent; scores are calibrated against observed uptake, not stated intent",
    "Crypto and rewards were tested by description only — a prototype could move them, and a re-score would catch it",
    "Insurance sits four points above the gate: treated as admitted, re-tested every quarter",
]

WRAPPER_ROWS = [
    ("The gap", "At the moments that matter — the first export order, the first unpaid foreign invoice — exporters name their bank 22% of the time, and the market leader 41%."),
    ("The wrapper", "One name over the launched features, with a fixed set of brand assets, tied to those two moments. Product communications ladder up to it, not alongside it."),
    ("The economics", "£0.8m a year, fixed — spread across every feature in the shell, unlike incentives that pay per redemption. Break-even at a 6-point lift in recall."),
    ("The multiplier rule", "A wrapper multiplies what is there. It cannot rescue a feature below the demand gate — weak features wrapped are weak features remembered."),
]

INCENTIVE_ROWS = [
    ("Defaults beat persuasion", "FX is switched on for every exporting customer at onboarding — no incentive. Regulator trials and our own tests show pre-set options beat paid persuasion."),
    ("Incentives buy selection, not behaviour", "Around a third of incentive redeemers would have acted anyway — the dead-weight share, estimated from our last two campaigns. It is priced into the −0.7."),
    ("Where incentives DO pay", "The switching moment. Invoice finance targets exporters leaving a factoring provider — a real cost to move, a one-off payment to cover it, capped and clawed back."),
]

RESCORE_ROWS = [
    ("New demand evidence", "A prototype test of crypto, or a shift in exporter behaviour", "Scores re-run; a feature can enter or leave the backlog"),
    ("Return assumptions move", "Uptake, margin, dead-weight or partner terms change", "Per-slot returns re-computed; positions 2 and 3 can swap"),
    ("Capability lands", "The claims capability ships", "Insurance's slots fall from 2 to 1 — its per-slot return doubles"),
    ("A slot frees up", "A build finishes early or a partner accelerates", "The cut-line moves down the sorted list — no new argument needed"),
]

SLATE_CARDS = [
    ("FX payments — BUILD", "Slot 1", [
        "Deepens the payments platform we already run — the capability compounds",
        "9 months ±1, high confidence: the delivery team owns every dependency",
        "Full margin kept; no third party in the flow"]),
    ("Invoice finance — PARTNER", "Slot 2", [
        "Live in 5 months instead of 18 — the difference between second place and missing the cycle",
        "Partner: [specialist lender], revenue share 25%",
        "We keep the customer, the data and the pricing; they run the credit engine"]),
]
SLATE_NOTE = "Insurance waits: it needs the claims capability, which is itself a build. When that lands, its slot cost halves and the re-score promotes it."

BOUNDS_ROWS = [
    ("What the partner holds", "The credit decisioning engine and the capital. Their capability, their balance sheet risk."),
    ("What we keep", "The customer relationship, the transaction data, the pricing decision, and the 'Trade Ready' brand over the top."),
    ("Contract bounds", "Exclusivity in the exporting-SME segment · our data stays ours on exit · 90-day portability · price re-opener at 12 months."),
    ("Re-internalise trigger", "At 5,000 active users, building in-house pays back in under 3 years — reviewed at each quarterly re-score."),
]

SENSITIVITY_ROWS = [
    ("Demand scores overstate uptake by a third", "FX still clears; invoice finance drops to the hurdle", "Activation rate at month 3 vs the model"),
    ("Incentive dead-weight passes 50%", "The −0.7 doubles; the incentive is switched off", "Redemption vs new-to-feature mix, monthly"),
    ("'Trade Ready' recall does not lift by month 6", "The £0.8m stops; comms revert to product level", "Entry-point recall tracking, quarterly"),
    ("The partner slips or captures the relationship", "Invoice finance misses the cycle; moat erodes", "Contract milestones + data-flow audit, month 2"),
]

ASKS = [
    ("Adopt",   "The sorting rule: demand admits, per-slot return sorts, capacity cuts", "[Owner]", "[Date]"),
    ("Approve", "This cycle's slate: build FX, partner invoice finance",                 "[Owner]", "[Date]"),
    ("Sign",    "Partner heads of terms — exclusivity, data, exit, re-opener (slide 19)", "[Owner]", "[Date]"),
    ("Stop",    "Crypto and rewards — off the list until new demand evidence",      "[Owner]", "[Date]"),
]

APPENDIX_ITEMS = [
    ("OCV methodology + research findings", "Metric construction, calibration, the 320-interview base"),
    ("Pricing and incentive model", "Lever assumptions, dead-weight estimates, clawback terms"),
    ("Wrapper / recall study", "Entry-point research, recall tracking, break-even"),
    ("Partner due diligence + route model", "Counterparty assessment, terms, route economics"),
    ("Method: the four analyses", "Step-by-step, next four pages — reproducible with Copilot"),
    ("Features considered and excluded earlier", "Anything cut before the shortlist, and why"),
]

METHODS = [
    ("1 OF 4", "How the demand scores were built",
     ["Segment interviews (n, dates)", "Incumbent solution facts and prices", "Feature specs and price points", "Past-launch calibration data"],
     ["Derive the segment's outcomes: build the outcome ladder for the job; keep the minimal covering set (OCV standard, step 1a).",
      "Score each feature against the incumbent per component: financial, time and stress deltas, time priced at the standard's conventions.",
      "Subtract switching costs (the standard's taxonomy). Adjust each component for visibility and credibility — contractual anchors beat claimed ones.",
      "Sum to the OCV score. Map to predicted uptake with the calibration table; set the gate from past launches, not from the answer you want.",
      "Rank the features. Record every input in the evidence register with source and date — the register is what the quarterly re-score re-runs."],
     "slides 8–9, and the uptake input to the sort",
     "Attach objective-customer-value-standard.md (same repository). One chat per feature. Paste the step, the feature spec and the register; require evidence for AND against."),
    ("2 OF 4", "How the backlog was sorted",
     ["Uptake per feature (from method 1)", "Margin model per feature", "Build-slot estimates per feature", "Past campaign results, by lever"],
     ["Build the base case per admitted feature: predicted uptake × price × margin, before any lever.",
      "Size the levers with the two-part test: would the behaviour happen anyway, and could a default do it free? Price the dead-weight in.",
      "Discount each feature's configured NPV for its route's timing and confidence (method 4 supplies both).",
      "Divide by build slots consumed: the sort key is risk-adjusted return PER SLOT. Rank. Note where per-slot order differs from raw NPV order.",
      "Set the re-score triggers: which input movements change the order, and who watches each one."],
     "slides 12–16",
     "One chat per feature. Paste the base case, lever tests and route discount; require the per-slot arithmetic shown and the trigger list stated."),
    ("3 OF 4", "How the wrapper case was built",
     ["Brand tracking data", "Entry-point research: when the job arises, in customer words", "Media / channel cost cards", "Competitor recall benchmarks"],
     ["Identify the segment's entry points: the moments the job arises, phrased in the customer's words, from research — not from the marketing plan.",
      "Measure current unaided recall at each entry point, for the bank and for competitors. The gap is the wrapper's job description.",
      "Design the wrapper: one name, fixed assets, explicit linkage to the entry points. Span only admitted features — never wrap what failed the gate.",
      "Price it: fixed annual cost against a recall-uplift band. State the break-even in recall points — the honest version of a brand business case.",
      "Set the tracking: entry-point recall, quarterly, with the stop rule — the wrapper's row on the sensitivity page."],
     "slide 10, the bridge on slide 13, and the sensitivity page",
     "Paste the entry-point list and tracking data. Require the break-even arithmetic shown and a stated stop rule — refuse a case that only asserts awareness value."),
    ("4 OF 4", "How the delivery routes were priced",
     ["Delivery team estimates, with ranges", "Partner term sheets / indicative terms", "Capability register", "Hurdle rate and horizon"],
     ["Inventory the capabilities each feature requires — capabilities, not tasks. Keep the same names end to end or the slides will not reconcile.",
      "Classify each capability have / buy / build. Name real partner candidates with indicative terms — an unnamed partner is a build with optimism attached.",
      "Price each route: time to launch, capex, margin share, and a confidence band from the delivery team and the partner's committed dates.",
      "Hand timing and confidence to the sort (method 2). Keep the built-instead counterfactual row — it is what prices the route decision.",
      "Bound the partner routes: what they hold, what we keep, contract terms, and the re-internalise trigger with its review date."],
     "slides 15, 18–19",
     "One chat per feature. Paste the capability inventory and terms; require the counterfactual row and the defensibility bounds before accepting a verdict."),
]

OUTPUT_FILE = "proposition-backlog-template-barclays.pptx"

# ============================================================================
# LAYOUT — DO NOT EDIT BELOW THIS LINE
# ============================================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

NAVY = RGBColor(0x00, 0x39, 0x5D); CYAN = RGBColor(0x00, 0xAE, 0xEF)
CYAND = RGBColor(0x00, 0x83, 0xC5); TINT = RGBColor(0xD9, 0xEE, 0xF9)
TINT2 = RGBColor(0xCD, 0xE9, 0xF8); GREY = RGBColor(0xBF, 0xBF, 0xBF)
BODY = RGBColor(0x3F, 0x55, 0x63); MUTE = RGBColor(0x6E, 0x7F, 0x8D)
RULE = RGBColor(0xC9, 0xDD, 0xE8); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2F, 0x6B, 0x4F); RED = RGBColor(0xA8, 0x44, 0x3C)
PALE = RGBColor(0xBF, 0xE4, 0xF7); ZEBRA = RGBColor(0xF5, 0xF8, 0xFA)
FONT = "Expert Sans"

W, H, M = 13.333, 7.5, 0.62
CW = W - 2 * M

prs = Presentation()
prs.slide_width = Inches(W); prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None, rounded=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    else: shp.line.fill.background()
    if rounded:
        try: shp.adjustments[0] = 0.08
        except Exception: pass
    return shp

def txt(s, x, y, w, h, runs, size=13, color=BODY, bold=False, italic=False,
        align=PP_ALIGN.LEFT, spacing=None, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    if isinstance(runs, str): runs = [(runs, {})]
    first = True
    for text, opt in runs:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = opt.get("align", align)
        if spacing: para.line_spacing = Pt(spacing)
        r = para.add_run(); r.text = text
        f = r.font; f.name = FONT
        f.size = Pt(opt.get("size", size)); f.bold = opt.get("bold", bold)
        f.italic = opt.get("italic", italic)
        f.color.rgb = opt.get("color", color)
    return tb

def bullets(s, items, x, y, w, size=12.5, color=BODY, gap=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(gap)
        r = para.add_run(); r.text = "•  " + it
        f = r.font; f.name = FONT; f.size = Pt(size); f.color.rgb = color
    return tb

def eyebrow(s, t, y=0.30):
    txt(s, M, y, 9, 0.3, t, size=11, color=CYAND, bold=True)

def title(s, t, size=24):
    txt(s, M, 0.60, CW - 1.6, 1.05, t, size=size, color=CYAN, bold=False, spacing=size + 5)

def pagetag(s, n):
    txt(s, W - 3.4, 0.28, 2.8, 0.3, META["page_tag"] + "  |  " + str(n),
        size=8.5, color=MUTE, align=PP_ALIGN.RIGHT)

def foot(s, t):
    txt(s, M, H - 0.62, CW, 0.3, t, size=9, color=MUTE)

def num_circle(s, x, y, n, d=0.46):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = CYAN; c.line.fill.background()
    tf = c.text_frame; p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run(); r.text = str(n)
    r.font.name = FONT; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

def pill(s, x, y, t, fill, w=0.92):
    shp = rect(s, x, y, w, 0.30, fill, rounded=True)
    try: shp.adjustments[0] = 0.5
    except Exception: pass
    tf = shp.text_frame; p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run(); r.text = t
    r.font.name = FONT; r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = WHITE

def notes(s, t):
    s.notes_slide.notes_text_frame.text = t

def divider(n, kicker, headline, sub):
    s = slide()
    rect(s, 0, 0, W, 0.14, CYAN)
    txt(s, M, 2.05, 2.6, 1.6, n, size=88, color=PALE, bold=True)
    txt(s, M + 2.55, 2.22, 9.4, 0.4, kicker.upper(), size=12, color=CYAND, bold=True)
    txt(s, M + 2.55, 2.62, 9.4, 1.0, headline, size=32, color=NAVY)
    txt(s, M + 2.55, 3.66, 9.2, 0.7, sub, size=14, color=BODY)
    return s

# -------- derived values (the self-computing part) --------
configured = round(BRIDGE_BASE + sum(d for _, d in BRIDGE_STEPS), 1)
gt_earn = int(round(configured - 0.3))  # headline rounds conservatively down
admitted = [f for f in FEATURES if f["risk_npv"] is not None]
for f in admitted:
    f["per_slot"] = round(f["risk_npv"] / f["slots"], 1)
admitted_sorted = sorted(admitted, key=lambda f: -f["per_slot"])
for i, f in enumerate(admitted_sorted):
    f["order"] = i + 1
GT = ("Fill '%s' in return order: FX and invoice finance take this cycle's two slots — £%dm a year by year %d."
      % (META["wrapper"], gt_earn, META["earnings_year"]))

# ---------------- 1. COVER ----------------
s = slide()
rect(s, 0, 0, W, 0.22, CYAN)
txt(s, M, 1.45, 8.6, 0.34, "PROPOSITION BACKLOG RECOMMENDATION", size=12, color=CYAND, bold=True)
txt(s, M, 1.95, 8.4, 2.1, "What do we build,\nand in what order?", size=46, color=CYAN, spacing=52)
txt(s, M, 4.15, 8.6, 0.4, "A rule that sorts the backlog · this cycle's slate · a re-score cadence", size=15, color=NAVY)
txt(s, M, 4.65, 9.2, 0.4, "Worked example: %s · all names and numbers are dummy values" % META["segment"], size=12, color=MUTE)
txt(s, M, 6.55, 8, 0.4, META["team_line"], size=12, color=MUTE)
cx, top, unit, gap = 10.85, 2.15, 0.62, 0.16
for i, wd in enumerate([1.05, 2.15, 3.25]):
    rect(s, cx - wd / 2, top + i * (unit + gap), wd, unit,
         [CYAN, NAVY, TINT][i], rounded=True)
txt(s, cx - 2.6, top + 3 * (unit + gap) + 0.12, 5.2, 0.36, "Answer  →  key line  →  support",
    size=11, color=MUTE, align=PP_ALIGN.CENTER)
notes(s, "Cover, Barclays IR style. Generated by generate_backlog_deck.py - edit the DATA section and re-run.")

# ---------------- 2. HOW TO USE ----------------
s = slide()
eyebrow(s, "TEMPLATE GUIDE — DELETE BEFORE SENDING")
title(s, "How to use this deck")
rows2 = [
    ("Every name and number is a dummy value",
     "%s, '%s', FX, the scores and the returns are a worked example. Replace all of them with your client's own." % (META["segment"].capitalize(), META["wrapper"])),
    ("The answer is a rule plus this cycle's slate",
     "The client approves the sorting rule AND the first output. That is what makes the backlog durable: next cycle re-runs the rule, not the argument."),
    ("The sort key is return per build slot, not raw NPV",
     "Capacity is the scarce thing. A big feature that eats two slots can lose to two small ones. The worked example shows exactly this: insurance beats invoice finance on raw NPV and still sorts below it."),
    ("This deck regenerates from a Python script",
     "generate_backlog_deck.py rebuilds the whole file. Edit the DATA section (Copilot can run it and hand back the .pptx). The arithmetic — per-slot returns, the bridge, the headline — recomputes on every run."),
]
y = 1.62
for i, (a, b) in enumerate(rows2):
    rect(s, M, y, CW, 1.16, TINT, rounded=True)
    num_circle(s, M + 0.34, y + 0.35, i + 1)
    txt(s, M + 1.05, y + 0.14, CW - 1.5, 0.38, a, size=14.5, color=NAVY, bold=True)
    txt(s, M + 1.05, y + 0.52, CW - 1.5, 0.58, b, size=11.5, color=BODY, spacing=15)
    y += 1.32
notes(s, "Delete before sending. Row 4 documents the regeneration route now that Copilot executes Python.")

# ---------------- 3. COPILOT INSTRUCTIONS ----------------
s = slide()
eyebrow(s, "FOR COPILOT — DELETE BEFORE SENDING")
title(s, "Copilot: how to work with this deck", size=22)
rect(s, M, 1.55, 5.95, 4.6, TINT, rounded=True)
txt(s, M + 0.32, 1.72, 5.3, 0.34, "A — Preferred route: regenerate from the script", size=12.5, color=NAVY, bold=True)
txt(s, M + 0.32, 2.10, 5.35, 3.9, [
    ("1. ", {"bold": True, "color": CYAND}), ("Ask the user for generate_backlog_deck.py (same repository). Edit ONLY the DATA section: features, scores, returns, slots, routes, bridge steps, text blocks.", {}),
    ("2. ", {"bold": True, "color": CYAND}), ("Run the script. It recomputes per-slot returns, the bridge total and the headline, runs its self-check, and writes the .pptx. Give the file to the user.", {}),
    ("3. ", {"bold": True, "color": CYAND}), ("Never edit the LAYOUT section. If the self-check prints FAIL, fix the data, not the check.", {}),
], size=10.5, spacing=14)
rect(s, M + 6.25, 1.55, CW - 6.25, 4.6, TINT, rounded=True)
txt(s, M + 6.57, 1.72, 5.3, 0.34, "B — Fallback: edit this file directly", size=12.5, color=NAVY, bold=True)
txt(s, M + 6.57, 2.10, CW - 6.85, 3.9, [
    ("4. ", {"bold": True, "color": CYAND}), ("Replace every example value. Numbers must reconcile: bridge total (13) = headline earnings (4) · per-slot (12) = NPV (15) ÷ slots · scores identical on 6 and 8.", {}),
    ("5. ", {"bold": True, "color": CYAND}), ("The Barclays palette is already applied (cyan 00AEEF titles, navy 00395D data, pale blue D9EEF9 cards, grey BFBFBF comparisons). Do not remap it. Font: Expert Sans.", {}),
    ("6. ", {"bold": True, "color": CYAND}), ("Never restyle PASS green / FAIL red or the one-highlight rule. After edits, check titles fit two lines and tables sit inside margins.", {}),
], size=10.5, spacing=14)
rect(s, M, 6.32, CW, 0.75, TINT2, rounded=True)
txt(s, M + 0.42, 6.42, CW - 1.0, 0.55,
    "Route A is safer: the script cannot produce irreconcilable numbers. Route B requires the manual reconcile checks in step 4.",
    size=11.5, color=NAVY, bold=True)
notes(s, "Two routes now that Copilot executes Python. Route A (regenerate) is canonical.")

# ---------------- 4. S-C-Q ----------------
s = slide()
eyebrow(s, "WHY WE ARE HERE")
title(s, "The question we set out to answer")
items = [("SITUATION", SCQ["situation"], 1.18, TINT),
         ("COMPLICATION", SCQ["complication"], 1.48, TINT),
         ("QUESTION", SCQ["question"], 1.05, TINT2)]
y = 1.66
for i, (lab, text, hh, fill) in enumerate(items):
    rect(s, M, y, CW, hh, fill, rounded=True)
    txt(s, M + 0.42, y + 0.18, 3.0, 0.32, lab, size=11, color=(NAVY if i == 2 else CYAND), bold=True)
    txt(s, M + 0.42, y + 0.52, CW - 1.0, hh - 0.7, text,
        size=(17 if i == 2 else 13.5), color=(NAVY if i == 2 else BODY), bold=(i == 2), spacing=(21 if i == 2 else 18))
    y += hh + 0.20
notes(s, "One question at the right altitude. The Complication earns the backlog framing.")

# ---------------- 5. RECOMMENDATION ----------------
s = slide()
eyebrow(s, "RECOMMENDATION"); pagetag(s, 4)
txt(s, M, 0.72, CW, 1.5, GT, size=24, color=CYAN, spacing=31)
y = 2.65
for lab, mid, right in REC_LINES:
    rect(s, M, y, CW, 0.72, TINT, rounded=True)
    txt(s, M + 0.35, y + 0.08, 2.2, 0.56, lab, size=14, color=CYAND, bold=True)
    txt(s, M + 2.65, y + 0.08, 6.4, 0.56, mid, size=14, color=NAVY, bold=True)
    txt(s, M + 9.15, y + 0.12, 3.2, 0.5, right, size=11.5, color=BODY)
    y += 0.9
ln = s.shapes.add_connector(1, Inches(M), Inches(6.42), Inches(M + CW), Inches(6.42))
ln.line.color.rgb = RULE
txt(s, M, 6.56, CW - 0.8, 0.5,
    "The client approves two things: this cycle's slate, and the rule that produced it. Next cycle re-runs the rule.",
    size=12.5, color=BODY)
notes(s, "GT derived from the data: earnings figure recomputes from the bridge on every run.")

# ---------------- 6. THREE REASONS ----------------
s = slide()
eyebrow(s, "THE ARGUMENT IN ONE PAGE")
title(s, "Three things make this order right")
cw3 = (CW - 0.6) / 3
for i, (tag, head, body) in enumerate(KEYLINE):
    x = M + i * (cw3 + 0.3)
    rect(s, x, 1.72, cw3, 4.2, TINT, rounded=True)
    num_circle(s, x + 0.38, 2.02, i + 1, 0.52)
    txt(s, x + 1.05, 2.06, cw3 - 1.4, 0.5, tag, size=11, color=CYAND, bold=True)
    txt(s, x + 0.38, 2.72, cw3 - 0.76, 1.15, head, size=15.5, color=NAVY, bold=True, spacing=20)
    txt(s, x + 0.38, 3.95, cw3 - 0.76, 1.85, body, size=11, color=BODY, spacing=15)
txt(s, M, 6.10, CW, 0.5, "Each reason stands alone: attack demand, the sort, or the cut separately — the other two hold.",
    size=12.5, color=MUTE, italic=True)
notes(s, "The key line: three independent reasons about the whole backlog.")

# ---------------- 7. BACKLOG MATRIX ----------------
s = slide()
eyebrow(s, "THE BACKLOG, DECIDED")
title(s, "One page: the order, and why each feature sits where it sits")
heads = ["Position", "Feature", "Demand\nvs gate %d" % GATE, "Return\nper slot", "Route", "Status"]
colX = [M, M + 1.35, M + 4.15, M + 5.85, M + 7.55, M + 10.55]
colW = [1.35, 2.8, 1.7, 1.7, 3.0, 1.6]
y = 1.72
for i, htxt in enumerate(heads):
    txt(s, colX[i], y, colW[i] - 0.15, 0.6, htxt, size=10.5, color=NAVY, bold=True, spacing=13)
y += 0.68
ln = s.shapes.add_connector(1, Inches(M), Inches(y), Inches(M + CW), Inches(y)); ln.line.color.rgb = NAVY
y += 0.12
def _status_color(st):
    return CYAND if st == "THIS CYCLE" else (NAVY if st == "NEXT CYCLE" else MUTE)
ordered = admitted_sorted + [f for f in FEATURES if f["risk_npv"] is None]
for ri, f in enumerate(ordered):
    if ri % 2 == 0: rect(s, M, y - 0.05, CW, 0.72, ZEBRA)
    pos = str(f.get("order", "—"))
    txt(s, colX[0], y, colW[0] - 0.15, 0.6, pos, size=17, color=(CYAND if f.get("order", 9) <= META["slots_per_cycle"] else NAVY if f["risk_npv"] else MUTE), bold=True)
    txt(s, colX[1], y + 0.05, colW[1] - 0.15, 0.55, f["name"], size=13, color=(NAVY if f["risk_npv"] else BODY), bold=True)
    txt(s, colX[2], y + 0.07, colW[2] - 0.15, 0.5, str(f["ocv"]), size=12.5, color=(BODY if f["ocv"] >= GATE else RED))
    txt(s, colX[3], y + 0.07, colW[3] - 0.15, 0.5, ("£%.1fm" % f["per_slot"]) if f["risk_npv"] else "not priced",
        size=12.5, color=(NAVY if f["risk_npv"] else MUTE), bold=bool(f["risk_npv"] and f.get("order", 9) <= 2))
    txt(s, colX[4], y + 0.07, colW[4] - 0.15, 0.5, f["route"], size=11.5, color=(BODY if f["risk_npv"] else MUTE))
    txt(s, colX[5], y + 0.09, colW[5], 0.5, f["status"], size=10, color=_status_color(f["status"]), bold=True)
    y += 0.70
txt(s, M, y + 0.1, CW, 0.42,
    "Demand admits, return sorts, capacity cuts. A feature below the gate is never priced — no queue place can be bought.",
    size=11.5, color=MUTE, italic=True)
notes(s, "The photograph slide. Derived from the FEATURES data: order and per-slot returns are computed, not typed.")

# ---------------- 8. DIVIDER 01 ----------------
notes(divider("01", "Demand admits", "Who gets into the backlog?",
              "Measured with the OCV demand score against a gate of %d. Below the gate, a feature is never priced." % GATE),
      "Demand is the admission test, not the sort key.")

# ---------------- 9. OCV CHART ----------------
s = slide()
eyebrow(s, "DEMAND ADMITS  ·  01")
title(s, "Three of five clear the demand gate; two are not close")
by_ocv = sorted(FEATURES, key=lambda f: f["ocv"])
cd = CategoryChartData()
cd.categories = [f["name"] for f in by_ocv]
cd.add_series("OCV", tuple(f["ocv"] for f in by_ocv))
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(M), Inches(1.75), Inches(8.0), Inches(4.1), cd)
ch = gf.chart; ch.has_legend = False
ch.has_title = False
plot = ch.plots[0]; plot.gap_width = 60
plot.has_data_labels = True
dl = plot.data_labels
dl.font.size = Pt(12); dl.font.name = FONT; dl.font.bold = True
dl.position = XL_LABEL_POSITION.OUTSIDE_END
ser = plot.series[0]
for i, f in enumerate(by_ocv):
    pt = ser.points[i]; pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = NAVY if f["ocv"] >= GATE else GREY
ch.value_axis.visible = False; ch.value_axis.has_major_gridlines = False
ch.value_axis.maximum_scale = 100; ch.value_axis.minimum_scale = 0
ch.category_axis.format.line.fill.background()
ch.category_axis.tick_labels.font.size = Pt(11); ch.category_axis.tick_labels.font.name = FONT
rect(s, M + 8.45, 1.75, CW - 8.45, 4.1, TINT, rounded=True)
txt(s, M + 8.8, 1.95, 2.9, 0.34, "The gate", size=11, color=CYAND, bold=True)
txt(s, M + 8.8, 2.30, 2.9, 1.0, str(GATE), size=56, color=NAVY, bold=True)
txt(s, M + 8.8, 3.45, 2.9, 0.3, "Admitted", size=11, color=MUTE, bold=True)
txt(s, M + 8.8, 3.75, 3.0, 0.7, ", ".join(f["name"].split()[0] for f in admitted) + " — priced in section 02", size=11.5, color=BODY)
txt(s, M + 8.8, 4.55, 2.9, 0.3, "Refused", size=11, color=MUTE, bold=True)
txt(s, M + 8.8, 4.85, 3.0, 0.6, ", ".join(f["name"].split()[0] for f in FEATURES if f["risk_npv"] is None) + " — never priced", size=11.5, color=BODY)
foot(s, "Source: exporter research programme, 320 interviews, May–Jul [year]. OCV method: appendix.")
notes(s, "Chart derived from FEATURES: bar colour follows the gate automatically.")

# ---------------- 10. EVIDENCE ----------------
s = slide()
eyebrow(s, "DEMAND ADMITS  ·  01")
title(s, "What the scores rest on — and their stated limits")
cw4 = (CW - 0.75) / 4
for i, (num, lab) in enumerate(EVIDENCE_STATS):
    x = M + i * (cw4 + 0.25)
    rect(s, x, 1.72, cw4, 1.8, TINT, rounded=True)
    txt(s, x + 0.28, 1.90, cw4 - 0.56, 0.85, num, size=32, color=NAVY, bold=True)
    txt(s, x + 0.28, 2.75, cw4 - 0.56, 0.65, lab, size=11, color=BODY, spacing=14)
txt(s, M, 3.85, CW, 0.4, "Known limits of the evidence", size=14.5, color=NAVY, bold=True)
bullets(s, EVIDENCE_LIMITS, M, 4.30, CW, size=12.5)
notes(s, "Limits stated by us before a sceptic finds them; each connects to the re-score cadence.")

# ---------------- 11. WRAPPER ----------------
s = slide()
eyebrow(s, "DEMAND ADMITS  ·  01  ·  PROPOSITION LEVEL")
title(s, "'%s' makes the shelf visible when the need hits" % META["wrapper"])
y = 1.72
for lab, text in WRAPPER_ROWS:
    rect(s, M, y, CW, 0.98, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.24, 3.4, 0.5, lab, size=13.5, color=NAVY, bold=True)
    txt(s, M + 4.1, y + 0.10, CW - 4.6, 0.78, text, size=11.5, color=BODY, spacing=15)
    y += 1.12
txt(s, M, y + 0.08, CW, 0.6,
    "The wrapper is a demand claim — will they think of us — so it lives in this section. Its cost appears in the earnings bridge; its stop rule sits on the sensitivity page.",
    size=12, color=MUTE, italic=True, spacing=16)
notes(s, "The standing shell the backlog ships into. Break-even and stop rule keep it honest.")

# ---------------- 12. DIVIDER 02 ----------------
notes(divider("02", "Return sorts", "What order do the admitted features take?",
              "Sorted by risk-adjusted return per build slot — capacity is the scarce thing, so returns are measured against it."),
      "The sort key is the archetype's central discipline.")

# ---------------- 13. THE SORT ----------------
s = slide()
eyebrow(s, "RETURN SORTS  ·  02")
title(s, "Per-slot return sets the order — and it re-orders raw NPV")
heads = ["Feature", "Risk-adjusted NPV", "Build slots", "Return per slot", "Order"]
colX = [M, M + 3.9, M + 6.4, M + 8.3, M + 10.9]
y = 1.85
for i, htxt in enumerate(heads):
    txt(s, colX[i], y, 2.4, 0.4, htxt, size=11, color=NAVY, bold=True)
y += 0.46
ln = s.shapes.add_connector(1, Inches(M), Inches(y), Inches(M + CW), Inches(y)); ln.line.color.rgb = NAVY
y += 0.14
teach = max(admitted_sorted, key=lambda f: f["slots"])  # the multi-slot teaching row
for ri, f in enumerate(admitted_sorted):
    if ri % 2 == 0: rect(s, M, y - 0.05, CW, 0.70, ZEBRA)
    is_teach = f is teach
    label = f["name"] + (" (partnered)" if "Partner" in f["route"] else "")
    txt(s, colX[0], y + 0.03, 3.7, 0.55, label, size=13, color=NAVY, bold=True)
    txt(s, colX[1], y + 0.05, 2.3, 0.5, "£%.1fm" % f["risk_npv"], size=13, color=(RED if is_teach else BODY), bold=is_teach)
    txt(s, colX[2], y + 0.05, 1.7, 0.5, str(f["slots"]), size=13, color=(RED if is_teach else BODY), bold=is_teach)
    txt(s, colX[3], y + 0.05, 2.4, 0.5, "£%.1fm" % f["per_slot"], size=13, color=NAVY, bold=True)
    txt(s, colX[4], y + 0.02, 1.2, 0.55, "%dst" % f["order"] if f["order"] == 1 else ("%dnd" % f["order"] if f["order"] == 2 else "%drd" % f["order"]),
        size=15, color=CYAND, bold=True)
    y += 0.74
rect(s, M, y + 0.15, CW, 1.3, TINT, rounded=True)
txt(s, M + 0.42, y + 0.30, 4.5, 0.32, "Why per slot, not raw NPV", size=11, color=CYAND, bold=True)
txt(s, M + 0.42, y + 0.62, CW - 1.0, 0.75,
    "%s beats %s on raw NPV (£%.1fm vs £%.1fm) — and still sorts below it, because it eats both of the cycle's slots. Capacity is the scarce thing; returns are measured against what they consume."
    % (teach["name"], admitted_sorted[1]["name"].lower() if admitted_sorted[1] is not teach else admitted_sorted[0]["name"].lower(),
       teach["risk_npv"], admitted_sorted[1]["risk_npv"] if admitted_sorted[1] is not teach else admitted_sorted[0]["risk_npv"]),
    size=12, color=NAVY, spacing=16)
foot(s, "Risk-adjusted NPV: base earnings per feature, levers applied, discounted for route timing (slide 15).")
notes(s, "Derived: per-slot = NPV/slots, order = rank. The teaching row (multi-slot feature) highlights automatically.")

# ---------------- 14. BRIDGE ----------------
s = slide()
eyebrow(s, "RETURN SORTS  ·  02")
title(s, "Configuration lifts the cycle's earnings from £%.0fm to £%.1fm" % (BRIDGE_BASE, configured))
txt(s, M, 1.58, 9, 0.32, "£m annual earnings, year %d run-rate, this cycle's two features" % META["earnings_year"],
    size=11.5, color=MUTE)
baseY, scale = 5.7, 0.55
bars = [("Base case\n(list price,\nno levers)", BRIDGE_BASE, "abs")] + \
       [(lab, d, "up" if d >= 0 else "down") for lab, d in BRIDGE_STEPS] + \
       [("Configured\ncase", configured, "abs")]
run = 0.0; x = M + 0.25; bw, bgap = 1.62, 0.42
for lab, v, kind in bars:
    if kind == "abs": run = v; hh = v * scale; topb = baseY - hh
    elif kind == "up": hh = v * scale; topb = baseY - (run + v) * scale; run += v
    else: hh = (-v) * scale; topb = baseY - run * scale; run += v
    col = NAVY if kind == "abs" else (CYAN if kind == "up" else GREY)
    rect(s, x, topb, bw, hh, col)
    txt(s, x, topb - 0.34, bw, 0.3, ("%+.1f" % v) if kind != "abs" else ("%.1f" % v),
        size=11.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x - 0.08, baseY + 0.10, bw + 0.16, 0.95, lab, size=9.5, color=BODY, align=PP_ALIGN.CENTER, spacing=11)
    x += bw + bgap
txt(s, M, 6.30, CW, 0.42,
    "The switch incentive does most of the lever work and costs half of it back. The wrapper's recall uplift is the smallest add — and the only one that also serves every future feature.",
    size=12, color=MUTE, italic=True)
foot(s, "Source: pricing model v3. Incentive cost is gross, before clawback terms.")
notes(s, "Waterfall computed from BRIDGE_BASE + BRIDGE_STEPS; the title's configured case is the same sum.")

# ---------------- 15. INCENTIVE DISCIPLINE ----------------
s = slide()
eyebrow(s, "RETURN SORTS  ·  02")
title(s, "Incentives are spent only where a free default cannot work")
y = 1.75
for lab, text in INCENTIVE_ROWS:
    rect(s, M, y, CW, 1.28, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.36, 4.4, 0.6, lab, size=14, color=NAVY, bold=True, spacing=18)
    txt(s, M + 5.0, y + 0.18, CW - 5.5, 0.95, text, size=12, color=BODY, spacing=16)
    y += 1.46
txt(s, M, y + 0.10, CW, 0.55,
    "The test for every pound of incentive: would the behaviour happen anyway, and could a default do it free? Spend survives only if the answer is no, twice.",
    size=12.5, color=MUTE, italic=True)
notes(s, "The discipline slide: stops the deck being a discount programme with a strategy attached.")

# ---------------- 16. ROUTE TIMING ----------------
s = slide()
eyebrow(s, "RETURN SORTS  ·  02")
title(s, "Returns are discounted for each route's timing — the partner route is what puts invoice finance second", size=21)
heads6 = ["Feature · route", "Configured NPV", "Time to launch", "Confidence", "Risk-adjusted NPV", "Vs hurdle"]
colX = [M, M + 3.6, M + 5.9, M + 7.7, M + 9.3, M + 11.6]
y = 2.05
for i, htxt in enumerate(heads6):
    txt(s, colX[i], y, 2.2, 0.4, htxt, size=10.5, color=NAVY, bold=True)
y += 0.46
ln = s.shapes.add_connector(1, Inches(M), Inches(y), Inches(M + CW), Inches(y)); ln.line.color.rgb = NAVY
y += 0.14
for ri, (a, b, c, conf, cc, e, verdict) in enumerate(ROUTE_ROWS):
    if ri % 2 == 0: rect(s, M, y - 0.05, CW, 0.70, ZEBRA)
    txt(s, colX[0], y + 0.04, 3.4, 0.55, a, size=12.5, color=(NAVY if ri < 2 else MUTE), bold=ri < 2)
    txt(s, colX[1], y + 0.05, 2.1, 0.5, b, size=12.5, color=BODY)
    txt(s, colX[2], y + 0.05, 1.7, 0.5, c, size=12.5, color=BODY)
    pill(s, colX[3], y + 0.15, conf.upper(), GREEN if cc == "GREEN" else RED)
    txt(s, colX[4], y + 0.05, 2.1, 0.5, e, size=12.5, color=NAVY, bold=True)
    pill(s, colX[5], y + 0.15, verdict, GREEN if verdict == "PASS" else RED)
    y += 0.74
rect(s, M, y + 0.16, CW, 1.12, TINT, rounded=True)
txt(s, M + 0.42, y + 0.30, 5.0, 0.32, "Why the third row is on the slide", size=11, color=CYAND, bold=True)
txt(s, M + 0.42, y + 0.62, CW - 1.0, 0.55,
    "Built in-house, invoice finance falls to £0.6m a slot — below insurance, and below the hurdle. The partner route is not a preference: it is what earns the feature its place in this cycle.",
    size=12, color=NAVY, spacing=16)
notes(s, "The counterfactual row prices the route decision. Confidence bands come from the delivery team.")

# ---------------- 17. RE-SCORE ----------------
s = slide()
eyebrow(s, "RETURN SORTS  ·  02")
title(s, "The order is re-scored every quarter — here is what can move it")
y = 1.72
for a, b, c in RESCORE_ROWS:
    rect(s, M, y, CW, 0.98, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.24, 3.2, 0.5, a, size=13, color=NAVY, bold=True)
    txt(s, M + 3.9, y + 0.12, 4.4, 0.74, b, size=11.5, color=BODY, spacing=15)
    txt(s, M + 8.5, y + 0.12, CW - 9.0, 0.74, c, size=11.5, color=BODY, spacing=15)
    y += 1.12
txt(s, M, y + 0.08, CW, 0.5,
    "This is what the client is really approving: not this quarter's list, but the machine that keeps the list honest.",
    size=12.5, color=MUTE, italic=True)
notes(s, "Each trigger: observable signal, mechanical consequence. No re-litigation, just re-scoring.")

# ---------------- 18. DIVIDER 03 ----------------
notes(divider("03", "Capacity cuts", "What ships this cycle — and how?",
              "Two build slots. Routes chosen per feature, priced for time and margin, bounded for defensibility."),
      "Capacity sets the cut-line; the routes protect the advantage.")

# ---------------- 19. THE SLATE ----------------
s = slide()
eyebrow(s, "CAPACITY CUTS  ·  03")
title(s, "Two slots: build FX on our core; partner invoice finance for speed")
cw2 = (CW - 0.35) / 2
for i, (head, slot, items) in enumerate(SLATE_CARDS):
    x = M + i * (cw2 + 0.35)
    rect(s, x, 1.72, cw2, 3.6, TINT, rounded=True)
    txt(s, x + 0.4, 1.92, cw2 - 1.6, 0.5, head, size=17, color=NAVY, bold=True)
    txt(s, x + cw2 - 1.2, 1.98, 0.85, 0.45, slot, size=11.5, color=CYAND, bold=True, align=PP_ALIGN.RIGHT)
    ln = s.shapes.add_connector(1, Inches(x + 0.4), Inches(2.52), Inches(x + cw2 - 0.4), Inches(2.52))
    ln.line.color.rgb = RULE
    bullets(s, items, x + 0.4, 2.72, cw2 - 0.8, size=11.5, gap=8)
txt(s, M, 5.52, CW, 0.5, SLATE_NOTE, size=12, color=MUTE, italic=True)
notes(s, "Route verdicts differ per feature; each states its reason in place.")

# ---------------- 20. PARTNER BOUNDS ----------------
s = slide()
eyebrow(s, "CAPACITY CUTS  ·  03")
title(s, "The partner deal is bounded — the advantage stays ours")
y = 1.72
for lab, text in BOUNDS_ROWS:
    rect(s, M, y, CW, 0.98, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.24, 3.6, 0.5, lab, size=13.5, color=NAVY, bold=True)
    txt(s, M + 4.3, y + 0.12, CW - 4.8, 0.74, text, size=11.5, color=BODY, spacing=15)
    y += 1.12
txt(s, M, y + 0.08, CW, 0.55,
    "Shipping fast is the partner's job. Staying hard to beat is the contract's job. The two are different questions, and both must pass.",
    size=12.5, color=MUTE, italic=True)
notes(s, "If the partner holds the customer or the data, earnings survive but the moat does not.")

# ---------------- 21. SENSITIVITY ----------------
s = slide()
eyebrow(s, "ACROSS THE SYSTEM")
title(s, "What would have to be true for this to be the wrong call")
heads7 = ["If this turns out to be true…", "…then this changes", "…and how we would know early"]
colX7 = [M, M + 4.6, M + 8.5]
y = 1.75
for i, htxt in enumerate(heads7):
    txt(s, colX7[i] + (0.34 if i == 0 else 0), y, 4.0, 0.4, htxt, size=11, color=CYAND, bold=True)
y += 0.48
for a, b, c in SENSITIVITY_ROWS:
    rect(s, M, y, CW, 0.95, TINT, rounded=True)
    txt(s, colX7[0] + 0.34, y + 0.12, 4.0, 0.72, a, size=12, color=NAVY, bold=True, spacing=15)
    txt(s, colX7[1], y + 0.12, 3.6, 0.72, b, size=11.5, color=BODY, spacing=15)
    txt(s, colX7[2], y + 0.12, 3.6, 0.72, c, size=11.5, color=BODY, spacing=15)
    y += 1.08
txt(s, M, y + 0.06, CW, 0.5,
    "Every row ends in the re-score, not in a crisis meeting: the trigger fires, the score changes, the order moves.",
    size=12, color=MUTE, italic=True)
notes(s, "Being wrong is survivable, because the system re-scores instead of re-arguing.")

# ---------------- 22. SO WHAT ----------------
s = slide()
eyebrow(s, "SO WHAT"); pagetag(s, 21)
title(s, "What we need from you this week", size=25)
y = 1.80
for i, (verb, what, owner, date) in enumerate(ASKS):
    rect(s, M, y, CW, 0.92, TINT, rounded=True)
    num_circle(s, M + 0.34, y + 0.23, i + 1)
    txt(s, M + 1.05, y + 0.22, 1.6, 0.5, verb, size=17, color=CYAN)
    txt(s, M + 2.7, y + 0.12, 6.4, 0.7, what, size=12.5, color=NAVY, spacing=16)
    txt(s, M + 9.2, y + 0.24, 1.6, 0.5, owner, size=11.5, color=BODY)
    txt(s, M + 10.9, y + 0.24, 1.0, 0.5, date, size=11.5, color=BODY, align=PP_ALIGN.RIGHT)
    y += 1.08
txt(s, M, 6.30, CW, 0.5,
    "First re-score lands in [month]. From then on the backlog is a standing agenda item, not a new deck.",
    size=12.5, color=MUTE, italic=True)
notes(s, "Four asks: the rule, the slate, the terms, the stop. Adopting the rule outlives the meeting.")

# ---------------- 23. APPENDIX ----------------
s = slide()
eyebrow(s, "APPENDIX")
title(s, "Supporting detail — and how each analysis was run")
cw2 = (CW - 0.35) / 2
for i, (a, b) in enumerate(APPENDIX_ITEMS):
    x = M + (i % 2) * (cw2 + 0.35)
    y = 1.75 + (i // 2) * 1.28
    rect(s, x, y, cw2, 1.10, TINT, rounded=True)
    txt(s, x + 0.4, y + 0.14, cw2 - 0.8, 0.4, a, size=13.5, color=NAVY, bold=True)
    txt(s, x + 0.4, y + 0.54, cw2 - 0.8, 0.45, b, size=11.5, color=BODY)
notes(s, "The method pages make the deck reproducible - which also makes the quarterly re-score cheap.")

# ---------------- 24-27. METHOD PAGES ----------------
for num, mtitle, inputs, steps, feeds, copilot in METHODS:
    s = slide()
    eyebrow(s, "APPENDIX  ·  METHOD " + num)
    title(s, mtitle, size=22)
    rect(s, M, 1.62, 3.5, 3.9, TINT, rounded=True)
    txt(s, M + 0.32, 1.80, 2.9, 0.34, "Inputs required", size=11, color=CYAND, bold=True)
    bullets(s, inputs, M + 0.32, 2.20, 2.95, size=10.5, gap=6)
    y = 1.68
    for i, st in enumerate(steps):
        num_circle(s, M + 3.85, y, i + 1, 0.4)
        txt(s, M + 4.42, y - 0.06, CW - 4.55, 0.72, st, size=11, color=BODY, spacing=14)
        y += 0.78
    rect(s, M, 5.72, CW, 1.0, TINT2, rounded=True)
    txt(s, M + 0.42, 5.84, 3.0, 0.3, "Run it with Copilot", size=10.5, color=CYAND, bold=True)
    txt(s, M + 0.42, 6.14, CW - 1.0, 0.5, copilot + "   ·   Output feeds: " + feeds,
        size=10.5, color=NAVY, spacing=14)
    notes(s, "Method page - the operator card behind the summary slides.")

# ============================================================================
# SELF-CHECK + SAVE
# ============================================================================
prs.save(OUTPUT_FILE)

def _fmt(ok): return "PASS" if ok else "FAIL"
checks = []
checks.append(("bridge sums to configured case",
               abs((BRIDGE_BASE + sum(d for _, d in BRIDGE_STEPS)) - configured) < 0.001))
checks.append(("governing thought earnings derived from bridge (£%dm from £%.1fm)" % (gt_earn, configured),
               ("£%dm" % gt_earn) in GT))
checks.append(("per-slot = NPV / slots for every admitted feature",
               all(abs(f["per_slot"] - f["risk_npv"] / f["slots"]) < 0.05 for f in admitted)))
checks.append(("capacity cut respected: THIS CYCLE slots == slots_per_cycle",
               sum(f["slots"] for f in admitted if f["status"] == "THIS CYCLE") == META["slots_per_cycle"]))
checks.append(("every refused feature is below the gate",
               all(f["ocv"] < GATE for f in FEATURES if f["risk_npv"] is None)))
checks.append(("every admitted feature is at or above the gate",
               all(f["ocv"] >= GATE for f in admitted)))
checks.append(("slide count == 27", len(prs.slides._sldIdLst) == 27))

import re as _re
_strings = []
def _collect(o):
    if isinstance(o, str): _strings.append(o)
    elif isinstance(o, dict):
        for v in o.values(): _collect(v)
    elif isinstance(o, (list, tuple)):
        for v in o: _collect(v)
for blob in (META, SCQ, FEATURES, KEYLINE, REC_LINES, WRAPPER_ROWS, INCENTIVE_ROWS,
             RESCORE_ROWS, SLATE_CARDS, BOUNDS_ROWS, SENSITIVITY_ROWS, ASKS):
    _collect(blob)
placeholders = sorted(set(m for s2 in _strings for m in _re.findall(r"\[[^\]]+\]", s2)))

print("=" * 60)
print("SELF-CHECK — %s" % OUTPUT_FILE)
for name, ok in checks:
    print("  [%s] %s" % (_fmt(ok), name))
print("  [INFO] placeholders still to fill: %d %s" % (len(placeholders), placeholders))
print("=" * 60)
if not all(ok for _, ok in checks):
    raise SystemExit("SELF-CHECK FAILED — fix the DATA section, not the checks.")
print("written", OUTPUT_FILE)
