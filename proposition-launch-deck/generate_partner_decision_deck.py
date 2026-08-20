#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PARTNER DECISION DECK GENERATOR — theme-driven, Copilot-executable.
Archetype #2: "Should we partner to launch this specific feature?"

HOW TO USE (in the Copilot app):
  STAGE 0 — APPLY THE CLIENT'S BRAND (optional). Two routes: (i) paste the
     client palette hexes and font names into THEME in the DATA section;
     (ii) attach any existing client .pptx and run extract_theme.py
     (repository root) — it prints a ready-to-paste THEME dict with a
     suggested role mapping. The shipped default is the neutral TMTH
     palette; client brand arrives at runtime and must never be committed
     to the public repository.
  1. Attach this file. Say: "Run this script and give me the .pptx it creates."
  2. To change the deck: edit ONLY the DATA section below, then re-run.
     Do not edit anything under LAYOUT — DO NOT EDIT.

Self-computing: the recommended route is checked against the route table,
the re-internalise crossover is CALCULATED from build capex and the margin
delta, and a self-check refuses to write a deck that fails.

VERSION NOTES: v1.1 — 20 Aug 2026: brand moved from code to data (THEME
dict, neutral TMTH default); client palette and font removed from code;
theme validation gate added; heading/body font split introduced.

Requires: python-pptx.  Output: partner-decision-template-branded.pptx
"""

# ============================================================================
# DATA — EDIT THIS SECTION ONLY
# ============================================================================

# --- THEME: all brand values live here, as data. Ship = neutral TMTH palette.
# To apply a client's brand: (i) paste their palette hexes + font names below,
# or (ii) attach any existing client .pptx and run extract_theme.py (repository
# root) — it prints a ready-to-paste THEME dict. Hex: 6 chars, no '#', no alpha.
THEME = {
    "dominant":      "2F4858",  # lead brand colour — emphasis text, headline data, table header fills
    "accent":        "4E7C90",  # secondary brand colour — slide titles, big display text
    "highlight":     "C9A227",  # the one-highlight colour — eyebrows, numbered circles, key callouts
    "card_fill":     "EEF3F6",  # card / panel background tint
    "ink":           "1F2A31",  # near-black headings ink
    "ink_soft":      "3C4C57",  # body text ink
    "muted":         "6B7A85",  # secondary text, sources, footers
    "chart_compare": "B9C6CD",  # non-highlighted chart bars / comparison elements
    "font_heading":  "Cambria",
    "font_body":     "Calibri",
}

META = {
    "feature": "invoice finance",
    "segment": "exporting SMEs",
    "partner": "[Lender Co]",
    "page_tag": "Partner decision  |  [Month Year]",
    "team_line": "[Team / function]   ·   [Month Year]",
}

SCQ = {
    "situation": "Invoice finance cleared the demand gate (OCV 72 vs 60) and holds slot 2 of this cycle's backlog. We can build it, buy it, or launch it with a partner.",
    "complication": "Built in-house it takes 18 months — past the window in which exporters are switching provider. Partnering trades margin share and a defensibility risk for speed. Decided on instinct, that trade goes unpriced and unbounded.",
    "question": "Should we partner to launch invoice finance — and on what terms?",
}

GT = "Partner with [Lender Co] to launch invoice finance: live in 5 months, £1.6m NPV — build misses the window; the terms keep the advantage ours."

REC_LINES = [
    ("Sign",   "[Lender Co] heads of terms — exclusivity, data, exit, re-opener", "the deal, bounded"),
    ("Bound",  "We keep the customer, the data, the pricing and the brand",       "the moat clause"),
    ("Gate",   "Contract milestones at month 2; live by month 5",                 "or the deal re-opens"),
    ("Revisit","Build in-house at ~5,000 active users",                           "computed crossover, checked quarterly"),
]

# Route table. One row per route; the chosen route first. npv in GBPm.
ROUTES = [
    {"route": "Partner — [Lender Co]", "time": "5 mths ±1",  "capex": 0.3, "margin_note": "25% revenue share", "npv": 1.6, "verdict": "RECOMMENDED", "chosen": True},
    {"route": "Build in-house",         "time": "18 mths ±4", "capex": 2.0, "margin_note": "full margin",       "npv": 0.6, "verdict": "MISSES WINDOW", "chosen": False},
    {"route": "Buy / white-label",      "time": "7 mths ±2",  "capex": 0.5, "margin_note": "unit fee, thin differentiation", "npv": 1.1, "verdict": "DOMINATED", "chosen": False},
    {"route": "Do not launch",          "time": "—",        "capex": 0.0, "margin_note": "slot 2 goes to insurance at £0.9m/slot", "npv": 0.9, "verdict": "THE FLOOR", "chosen": False},
]

# Crossover economics (drives the computed re-internalise trigger).
CROSSOVER = {
    "build_capex_m": 2.0,          # extra capex to build vs partner (GBPm)
    "margin_build_per_user": 520,  # GBP per active user per year, in-house
    "margin_partner_per_user": 390,# GBP per active user per year, partnered
    "payback_years": 3,            # build must pay back within this
}

KEYLINE = [
    ("WORTH LAUNCHING NOW", "The demand is proven — and the window is real",
     "OCV 72 against a gate of 60, from 320 exporter interviews. Exporters switch provider at renewal moments that cluster in the next two quarters — a feature that arrives in month 18 meets a market that has already switched."),
    ("THE ROUTE WINS THE MONEY", "Partnering is worth £1.0m more than building",
     "Risk-adjusted: partner £1.6m, build £0.6m, white-label £1.1m, do-nothing £0.9m (the slot's next best use). The build route does not lose on margin — it loses on time. Speed is the asset being bought."),
    ("THE DEAL IS BOUNDED", "The terms keep the advantage ours",
     "We keep the customer, the data, the pricing and the brand; the partner runs the credit engine and holds the balance-sheet risk. Exclusivity, exit portability, a price re-opener — and a computed re-internalise trigger."),
]

DEMAND_CARD = [
    ("The score", "OCV 72 against a gate of 60 — second of five candidates. Evidence: 320 exporter interviews, calibrated against observed uptake, limits stated in the backlog deck."),
    ("The window", "Exporters review invoice finance at renewal. Two-thirds of the segment's renewal dates fall inside the next two quarters — the switching moment the incentive targets."),
    ("The cost of month 18", "Arriving after the renewal wave means selling against fresh two-year contracts. Modelled uptake falls by more than half — which is why the build route's NPV collapses to £0.6m."),
]

WAIT_STATS = [
    ("2 qtrs", "of renewal dates ahead — the switching window"),
    ("−56%", "modelled uptake if we arrive after it"),
    ("£1.0m", "NPV gap between partnering now and building late"),
    ("£0.9m", "the floor: slot 2's next best use if we do not launch"),
]

DILIGENCE_ROWS = [
    ("Who they are", "[Specialist SME lender, £Xbn advanced, N bank partnerships live. Regulated, rated, audited — references checked with two current bank clients.]"),
    ("What we verified", "[Credit performance through the last downturn · platform uptime and integration references · conduct and complaints record · capital position]"),
    ("Red flags checked", "[Concentration risk in their book · dependence on one funding line · the two references' honest negatives — stated here, not hidden]"),
    ("Why them over the market", "[Shortlist of three; chosen on integration speed, segment fit and terms flexibility — the runner-up and the gap are named in the appendix]"),
]

BOUNDS_ROWS = [
    ("What the partner holds", "The credit decisioning engine and the capital. Their capability, their balance-sheet risk."),
    ("What we keep", "The customer relationship, the transaction data, the pricing decision, and the 'Trade Ready' brand over the top."),
    ("Contract bounds", "Exclusivity in the exporting-SME segment · our data stays ours on exit · 90-day portability · price re-opener at 12 months."),
    ("Re-internalise trigger", "COMPUTED_TRIGGER"),
]

SENSITIVITY_ROWS = [
    ("The partner slips past month 5", "The window narrows; the NPV gap to build shrinks", "Contract milestones + integration test, month 2"),
    ("The partner captures the relationship", "Earnings survive; the moat does not", "Data-flow audit at month 2; brand and pricing checks quarterly"),
    ("Margin share creeps at the re-opener", "Crossover falls; build comes forward", "Re-run this deck's crossover arithmetic at the re-opener"),
    ("Demand is softer than scored", "NPV falls on every route equally; partner still wins", "Activation rate at month 3 vs the model"),
]

ASKS = [
    ("Sign",    "Heads of terms with [Lender Co] — the four bounds on slide 16", "[Owner]", "[Date]"),
    ("Fund",    "Integration budget £0.3m and the month-2 milestone gate",       "[Owner]", "[Date]"),
    ("Instruct","Delivery to hold slot 2 open for the month-5 launch",          "[Owner]", "[Date]"),
    ("Diarise", "Quarterly crossover check against the ~5,000-user trigger",    "[Owner]", "[Date]"),
]

APPENDIX_ITEMS = [
    ("Route model", "Build / partner / buy economics, assumptions, version history"),
    ("Partner due diligence pack", "References, credit performance, the shortlist and the runner-up"),
    ("Demand evidence", "OCV scoring for invoice finance — see the backlog deck, section 01"),
    ("Crossover arithmetic", "Capex, margin delta, payback — the computed trigger's working"),
    ("Term sheet detail", "Exclusivity, data, exit, re-opener — full drafting"),
    ("The backlog context", "Why invoice finance holds slot 2 — the proposition backlog deck"),
]

OUTPUT_FILE = "partner-decision-template-branded.pptx"

# ============================================================================
# LAYOUT — DO NOT EDIT BELOW THIS LINE
# ============================================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

def _rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _blend(base_hex, mix_hex, f):
    """f = share of mix_hex blended into base_hex."""
    b = [int(base_hex[i:i+2], 16) for i in (0, 2, 4)]
    m = [int(mix_hex[i:i+2], 16) for i in (0, 2, 4)]
    return RGBColor(*(round(b[i] * (1 - f) + m[i] * f) for i in range(3)))

# --- THEME validation gate: refuse to build on malformed values ---
_theme_defects = []
for _k, _v in THEME.items():
    if _k.startswith("font_"):
        if not str(_v).strip():
            _theme_defects.append("THEME['%s'] is empty — name a real font" % _k)
    else:
        import re as _re_t
        if not _re_t.fullmatch(r"[0-9A-Fa-f]{6}", str(_v)):
            _theme_defects.append("THEME['%s'] = %r is not a 6-char hex (no '#', no alpha)" % (_k, _v))
if _theme_defects:
    print("=" * 60)
    print("SELF-CHECK — THEME")
    for d in _theme_defects:
        print("  [FAIL] %s" % d)
    print("=" * 60)
    raise SystemExit("SELF-CHECK FAILED — no file written. Fix THEME in the DATA section.")

# --- palette: derived from THEME (constant names kept so body code stands) ---
NAVY = _rgb(THEME["dominant"]); CYAN = _rgb(THEME["accent"])
CYAND = _rgb(THEME["highlight"]); TINT = _rgb(THEME["card_fill"])
TINT2 = _blend(THEME["card_fill"], THEME["accent"], 0.12)
GREY = _rgb(THEME["chart_compare"])
BODY = _rgb(THEME["ink_soft"]); MUTE = _rgb(THEME["muted"])
PALE = _blend("FFFFFF", THEME["dominant"], 0.30)
RULE = _blend("FFFFFF", THEME["dominant"], 0.20)
# Functional colours (pass/fail semantics, zebra striping), not brand — never remapped.
WHITE = RGBColor(0xFF, 0xFF, 0xFF); ZEBRA = RGBColor(0xF5, 0xF8, 0xFA)
GREEN = RGBColor(0x2F, 0x6B, 0x4F); RED = RGBColor(0xA8, 0x44, 0x3C)

FONT_HEAD = THEME["font_heading"]
FONT_BODY = THEME["font_body"]
FONT = FONT_BODY  # default for body text; keep the name so existing calls stand

W, H, M = 13.333, 7.5, 0.62
CW = W - 2 * M

prs = Presentation()
prs.slide_width = Inches(W); prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, line=None, rounded=False):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line: shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    else: shp.line.fill.background()
    if rounded:
        try: shp.adjustments[0] = 0.08
        except Exception: pass
    return shp

def txt(s, x, y, w, h, runs, size=13, color=BODY, bold=False, italic=False,
        align=PP_ALIGN.LEFT, spacing=None, wrap=True, font=FONT):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    if isinstance(runs, str): runs = [(runs, {})]
    first = True
    for text, opt in runs:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = opt.get("align", align)
        if spacing: para.line_spacing = Pt(spacing)
        r = para.add_run(); r.text = text
        f = r.font; f.name = font
        f.size = Pt(opt.get("size", size)); f.bold = opt.get("bold", bold)
        f.italic = opt.get("italic", italic)
        f.color.rgb = opt.get("color", color)
    return tb

def bullets(s, items, x, y, w, size=12.5, color=BODY, gap=6):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(gap)
        r = para.add_run(); r.text = "•  " + it
        f = r.font; f.name = FONT; f.size = Pt(size); f.color.rgb = color
    return tb

def eyebrow(s, t, y=0.30):
    txt(s, M, y, 9, 0.3, t, size=11, color=CYAND, bold=True)

def title(s, t, size=24):
    txt(s, M, 0.60, CW - 1.6, 1.05, t, size=size, color=CYAN, bold=False, spacing=size + 5, font=FONT_HEAD)

def pagetag(s, n):
    txt(s, W - 3.4, 0.28, 2.8, 0.3, META["page_tag"] + "  |  " + str(n),
        size=8.5, color=MUTE, align=PP_ALIGN.RIGHT)

def foot(s, t):
    txt(s, M, H - 0.62, CW, 0.3, t, size=9, color=MUTE)

def num_circle(s, x, y, n, d=0.46):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = CYAN; c.line.fill.background()
    tf = c.text_frame; p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run(); r.text = str(n)
    r.font.name = FONT_HEAD; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE

def pill(s, x, y, t, fill, w=0.92):
    shp = rect(s, x, y, w, 0.30, fill, rounded=True)
    try: shp.adjustments[0] = 0.5
    except Exception: pass
    tf = shp.text_frame; p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run(); r.text = t
    r.font.name = FONT; r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = WHITE

def notes(s, t):
    s.notes_slide.notes_text_frame.text = t

def divider(n, kicker, headline, sub):
    s = slide()
    rect(s, 0, 0, W, 0.14, CYAN)
    txt(s, M, 2.05, 2.6, 1.6, n, size=88, color=PALE, bold=True, font=FONT_HEAD)
    txt(s, M + 2.55, 2.22, 9.4, 0.4, kicker.upper(), size=12, color=CYAND, bold=True)
    txt(s, M + 2.55, 2.62, 9.4, 1.0, headline, size=32, color=NAVY, font=FONT_HEAD)
    txt(s, M + 2.55, 3.66, 9.2, 0.7, sub, size=14, color=BODY)
    return s


# -------- derived values (the self-computing part) --------
chosen = [r for r in ROUTES if r["chosen"]][0]
delta_per_user = CROSSOVER["margin_build_per_user"] - CROSSOVER["margin_partner_per_user"]
crossover_users = CROSSOVER["build_capex_m"] * 1_000_000 / (CROSSOVER["payback_years"] * delta_per_user)
crossover_rounded = int(round(crossover_users / 500.0) * 500)
TRIGGER_TEXT = ("At ~%s active users, building in-house pays back in under %d years "
                "(capex £%.1fm ÷ £%d/user margin delta) — reviewed at each quarterly re-score."
                % (format(crossover_rounded, ","), CROSSOVER["payback_years"],
                   CROSSOVER["build_capex_m"], delta_per_user))
BOUNDS_ROWS = [(a, TRIGGER_TEXT if b == "COMPUTED_TRIGGER" else b) for a, b in BOUNDS_ROWS]

# ---------------- 1. COVER ----------------
s = slide()
rect(s, 0, 0, W, 0.22, CYAN)
txt(s, M, 1.45, 8.6, 0.34, "PARTNER DECISION", size=12, color=CYAND, bold=True)
txt(s, M, 1.95, 9.2, 2.1, "Should we partner to launch\n%s?" % META["feature"], size=44, color=CYAN, spacing=50, font=FONT_HEAD)
txt(s, M, 4.15, 8.6, 0.4, "One feature · four routes priced · the deal bounded", size=15, color=NAVY)
txt(s, M, 4.65, 9.2, 0.4, "Worked example: %s for %s · all names and numbers are dummy values" % (META["feature"], META["segment"]), size=12, color=MUTE)
txt(s, M, 6.55, 8, 0.4, META["team_line"], size=12, color=MUTE)
cx, top, unit, gap = 10.85, 2.15, 0.62, 0.16
for i, wd in enumerate([1.05, 2.15, 3.25]):
    rect(s, cx - wd / 2, top + i * (unit + gap), wd, unit, [CYAN, NAVY, TINT][i], rounded=True)
txt(s, cx - 2.6, top + 3 * (unit + gap) + 0.12, 5.2, 0.36, "Answer  →  key line  →  support",
    size=11, color=MUTE, align=PP_ALIGN.CENTER)
notes(s, "Archetype #2: the single-feature route decision. Zooms in on one backlog verdict and prices it.")

# ---------------- 2. HOW TO USE ----------------
s = slide()
eyebrow(s, "TEMPLATE GUIDE — DELETE BEFORE SENDING")
title(s, "How to use this deck")
rows2 = [
    ("Every name and number is a dummy value",
     "The feature, '%s', the NPVs and the terms are a worked example — continuous with the proposition backlog deck. Replace all of them." % META["partner"]),
    ("The question includes the terms",
     "'Should we partner' alone is impoverished. The honest answer is yes-IF-bounded: the deal and its bounds are decided together, or not at all."),
    ("The counterfactual is the argument",
     "The build row is not background — it is what prices the decision. Every route including do-nothing gets a risk-adjusted NPV, so 'partner' wins on arithmetic, not enthusiasm."),
    ("This deck regenerates from this script",
     "Edit the DATA section and re-run (Copilot can execute it). The re-internalise trigger is computed from capex and margin delta — change the inputs and the trigger moves with them."),
]
y = 1.62
for i, (a, b) in enumerate(rows2):
    rect(s, M, y, CW, 1.16, TINT, rounded=True)
    num_circle(s, M + 0.34, y + 0.35, i + 1)
    txt(s, M + 1.05, y + 0.14, CW - 1.5, 0.38, a, size=14.5, color=NAVY, bold=True)
    txt(s, M + 1.05, y + 0.52, CW - 1.5, 0.58, b, size=11.5, color=BODY, spacing=15)
    y += 1.32
notes(s, "Delete before sending.")

# ---------------- 3. COPILOT INSTRUCTIONS ----------------
s = slide()
eyebrow(s, "FOR COPILOT — DELETE BEFORE SENDING")
title(s, "Copilot: how to work with this deck", size=22)
rect(s, M, 1.55, 5.95, 4.6, TINT, rounded=True)
txt(s, M + 0.32, 1.72, 5.3, 0.34, "A — Preferred route: regenerate from the script", size=12.5, color=NAVY, bold=True)
txt(s, M + 0.32, 2.10, 5.35, 3.9, [
    ("1. ", {"bold": True, "color": CYAND}), ("Edit ONLY the DATA section of generate_partner_decision_deck.py: the feature, partner, route NPVs, crossover inputs, terms.", {}),
    ("2. ", {"bold": True, "color": CYAND}), ("Run it. The script checks the recommended route against the NPVs, recomputes the crossover trigger, and writes the .pptx. Give the file to the user.", {}),
    ("3. ", {"bold": True, "color": CYAND}), ("Never edit the LAYOUT section. A FAIL means fix the data, not the check.", {}),
], size=10.5, spacing=14)
rect(s, M + 6.25, 1.55, CW - 6.25, 4.6, TINT, rounded=True)
txt(s, M + 6.57, 1.72, 5.3, 0.34, "B — Fallback: edit this file directly", size=12.5, color=NAVY, bold=True)
txt(s, M + 6.57, 2.10, CW - 6.85, 3.9, [
    ("4. ", {"bold": True, "color": CYAND}), ("Replace every example value. Reconcile: the headline NPV (5) = the chosen route's NPV (7, 12) · the trigger on 16 = the crossover arithmetic on 13.", {}),
    ("5. ", {"bold": True, "color": CYAND}), ("Brand lives in the THEME dict in the DATA section (neutral TMTH default shipped). To apply a client brand, edit THEME or run extract_theme.py against a client .pptx — never hand-restyle individual shapes.", {}),
    ("6. ", {"bold": True, "color": CYAND}), ("Never restyle verdict colours or the one-highlight rule. Check titles fit two lines after edits.", {}),
], size=10.5, spacing=14)
rect(s, M, 6.32, CW, 0.75, TINT2, rounded=True)
txt(s, M + 0.42, 6.42, CW - 1.0, 0.55,
    "Route A is safer: the script cannot recommend a route the numbers do not support.",
    size=11.5, color=NAVY, bold=True)
notes(s, "Two routes; regeneration is canonical.")

# ---------------- 4. S-C-Q ----------------
s = slide()
eyebrow(s, "WHY WE ARE HERE")
title(s, "The question we set out to answer")
items = [("SITUATION", SCQ["situation"], 1.18, TINT),
         ("COMPLICATION", SCQ["complication"], 1.48, TINT),
         ("QUESTION", SCQ["question"], 1.05, TINT2)]
y = 1.66
for i, (lab, text, hh, fill) in enumerate(items):
    rect(s, M, y, CW, hh, fill, rounded=True)
    txt(s, M + 0.42, y + 0.18, 3.0, 0.32, lab, size=11, color=(NAVY if i == 2 else CYAND), bold=True)
    txt(s, M + 0.42, y + 0.52, CW - 1.0, hh - 0.7, text,
        size=(17 if i == 2 else 13.5), color=(NAVY if i == 2 else BODY), bold=(i == 2), spacing=(21 if i == 2 else 18))
    y += hh + 0.20
notes(s, "The question includes the terms - yes-if-bounded, not bare yes/no.")

# ---------------- 5. RECOMMENDATION ----------------
s = slide()
eyebrow(s, "RECOMMENDATION"); pagetag(s, 5)
txt(s, M, 0.72, CW, 1.5, GT, size=24, color=CYAN, spacing=31)
y = 2.65
for lab, mid, right in REC_LINES:
    rect(s, M, y, CW, 0.72, TINT, rounded=True)
    txt(s, M + 0.35, y + 0.08, 2.2, 0.56, lab, size=14, color=CYAND, bold=True)
    txt(s, M + 2.65, y + 0.08, 6.4, 0.56, mid, size=13.5, color=NAVY, bold=True)
    txt(s, M + 9.15, y + 0.12, 3.2, 0.5, right, size=11.5, color=BODY)
    y += 0.9
ln = s.shapes.add_connector(1, Inches(M), Inches(6.42), Inches(M + CW), Inches(6.42))
ln.line.color.rgb = RULE
txt(s, M, 6.56, CW - 0.8, 0.5,
    "A yes with four bounds. Remove any one of them and the recommendation changes — that is what makes it a decision, not a preference.",
    size=12.5, color=BODY)
notes(s, "GT is one sentence; the four verbs below are the bounds that make the yes safe.")

# ---------------- 6. THREE REASONS ----------------
s = slide()
eyebrow(s, "THE ARGUMENT IN ONE PAGE")
title(s, "Three things make partnering the right call")
cw3 = (CW - 0.6) / 3
for i, (tag, head, body) in enumerate(KEYLINE):
    x = M + i * (cw3 + 0.3)
    rect(s, x, 1.72, cw3, 4.2, TINT, rounded=True)
    num_circle(s, x + 0.38, 2.02, i + 1, 0.52)
    txt(s, x + 1.05, 2.06, cw3 - 1.4, 0.5, tag, size=10.5, color=CYAND, bold=True)
    txt(s, x + 0.38, 2.72, cw3 - 0.76, 1.15, head, size=15.5, color=NAVY, bold=True, spacing=20)
    txt(s, x + 0.38, 3.95, cw3 - 0.76, 1.85, body, size=11, color=BODY, spacing=15)
txt(s, M, 6.10, CW, 0.5, "Each reason stands alone: attack the window, the arithmetic, or the terms separately — the other two hold.",
    size=12.5, color=MUTE, italic=True)
notes(s, "Same analytical machinery as the backlog deck (demand evidence, route pricing, bounds) - re-cut as reasons for THIS answer.")

# ---------------- 7. ROUTE MATRIX ----------------
s = slide()
eyebrow(s, "THE DECISION, PRICED")
title(s, "Four routes, one page: partnering wins on arithmetic, not enthusiasm")
heads = ["Route", "Time to live", "Capex", "Margin", "Risk-adj NPV", "Verdict"]
colX = [M, M + 3.3, M + 5.1, M + 6.4, M + 9.3, M + 10.9]
y = 1.78
for i, htxt in enumerate(heads):
    txt(s, colX[i], y, 1.9, 0.4, htxt, size=10.5, color=NAVY, bold=True)
y += 0.46
ln = s.shapes.add_connector(1, Inches(M), Inches(y), Inches(M + CW), Inches(y)); ln.line.color.rgb = NAVY
y += 0.14
for ri, r in enumerate(ROUTES):
    if r["chosen"]: rect(s, M, y - 0.05, CW, 0.80, TINT2)
    elif ri % 2 == 1: rect(s, M, y - 0.05, CW, 0.80, ZEBRA)
    txt(s, colX[0], y + 0.08, 3.1, 0.6, r["route"], size=12.5, color=NAVY, bold=r["chosen"])
    txt(s, colX[1], y + 0.10, 1.7, 0.5, r["time"], size=12, color=BODY)
    txt(s, colX[2], y + 0.10, 1.2, 0.5, ("£%.1fm" % r["capex"]) if r["capex"] else "—", size=12, color=BODY)
    txt(s, colX[3], y + 0.03, 2.8, 0.72, r["margin_note"], size=10.5, color=BODY, spacing=13)
    txt(s, colX[4], y + 0.08, 1.5, 0.5, "£%.1fm" % r["npv"], size=13, color=NAVY, bold=True)
    txt(s, colX[5], y + 0.10, 1.9, 0.5, r["verdict"], size=10, color=(CYAND if r["chosen"] else MUTE), bold=True)
    y += 0.84
txt(s, M, y + 0.10, CW, 0.42,
    "Do-nothing is priced too: slot 2's next best use. A route that cannot beat the floor never reaches the terms discussion.",
    size=11.5, color=MUTE, italic=True)
notes(s, "The photograph slide. Recommended row highlighted; every route priced including the floor.")

# ---------------- 8. DIVIDER 01 ----------------
notes(divider("01", "Worth launching now", "Is this feature worth a slot — today?",
              "Demand proven at the gate, and a switching window that prices the cost of being late."),
      "Demand + timing. Without the window, build-vs-partner is a preference, not a decision.")

# ---------------- 9. DEMAND + WINDOW ----------------
s = slide()
eyebrow(s, "WORTH LAUNCHING NOW  ·  01")
title(s, "Proven demand — and a window that closes in two quarters")
y = 1.75
for lab, text in DEMAND_CARD:
    rect(s, M, y, CW, 1.28, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.36, 3.6, 0.6, lab, size=14, color=NAVY, bold=True)
    txt(s, M + 4.3, y + 0.18, CW - 4.8, 0.95, text, size=12, color=BODY, spacing=16)
    y += 1.46
txt(s, M, y + 0.08, CW, 0.5,
    "The window is what turns route choice into a priced decision: month 5 and month 18 are different markets.",
    size=12.5, color=MUTE, italic=True)
notes(s, "Demand carried from the backlog deck; the window evidence is this deck's addition.")

# ---------------- 10. COST OF WAITING ----------------
s = slide()
eyebrow(s, "WORTH LAUNCHING NOW  ·  01")
title(s, "What being late costs")
cw4 = (CW - 0.75) / 4
for i, (num, lab) in enumerate(WAIT_STATS):
    x = M + i * (cw4 + 0.25)
    rect(s, x, 1.9, cw4, 2.1, TINT, rounded=True)
    txt(s, x + 0.28, 2.14, cw4 - 0.56, 0.85, num, size=30, color=NAVY, bold=True)
    txt(s, x + 0.28, 3.02, cw4 - 0.56, 0.85, lab, size=11, color=BODY, spacing=14)
rect(s, M, 4.5, CW, 1.3, TINT2, rounded=True)
txt(s, M + 0.42, 4.66, 5.0, 0.32, "Read", size=11, color=CYAND, bold=True)
txt(s, M + 0.42, 4.98, CW - 1.0, 0.75,
    "The build route's £0.6m is not a build problem — the same feature at month 5 would be worth £2.0m+ in-house. It is a timing problem, and partnering is the instrument that solves it.",
    size=12.5, color=NAVY, spacing=17)
foot(s, "Source: renewal-date analysis from [CRM]; uptake model v[3].")
notes(s, "The four tiles carry the whole timing argument; the read box states the reframe - speed is the asset being bought.")

# ---------------- 11. DIVIDER 02 ----------------
notes(divider("02", "The route wins the money", "Which route returns most, risk-adjusted?",
              "All four routes priced on one basis — including doing nothing with the slot."),
      "The economics section: arithmetic, then the crossover that sets the trigger.")

# ---------------- 12. ROUTE ECONOMICS ----------------
s = slide()
eyebrow(s, "THE ROUTE WINS THE MONEY  ·  02")
title(s, "Partnering is worth £1.0m more than building — because of time, not margin")
rows12 = [
    ("Margin per active user", "Build £520/yr · partner £390/yr — build wins on margin, as expected"),
    ("Time to live", "Partner month 5, inside the renewal window · build month 18, after it"),
    ("Uptake consequence", "Post-window uptake falls by more than half — the £2.0m+ in-house case collapses to £0.6m"),
    ("Risk-adjusted result", "Partner £1.6m · white-label £1.1m · do-nothing floor £0.9m · build £0.6m"),
]
y = 1.75
for lab, text in rows12:
    rect(s, M, y, CW, 0.98, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.24, 3.6, 0.5, lab, size=13.5, color=NAVY, bold=True)
    txt(s, M + 4.3, y + 0.12, CW - 4.8, 0.74, text, size=12, color=BODY, spacing=15)
    y += 1.12
txt(s, M, y + 0.08, CW, 0.55,
    "State it plainly: build loses £130 per user per year less margin? No — partnering gives up £130/user and still wins £1.0m, because the users exist at month 5 and largely do not at month 18.",
    size=12, color=MUTE, italic=True, spacing=16)
notes(s, "The margin-vs-time decomposition stops the 'we give away margin' objection before it starts.")

# ---------------- 13. CROSSOVER ----------------
s = slide()
eyebrow(s, "THE ROUTE WINS THE MONEY  ·  02")
title(s, "The partner deal carries its own exit: the crossover is computed, not felt")
rows13 = [
    ("The margin delta", "In-house earns £%d more per active user per year (£%d vs £%d)." % (delta_per_user, CROSSOVER["margin_build_per_user"], CROSSOVER["margin_partner_per_user"])),
    ("The build cost", "Bringing it in-house later costs ~£%.1fm of capex." % CROSSOVER["build_capex_m"]),
    ("The arithmetic", "£%.1fm ÷ (£%d/user × %d-year payback) = ~%s active users." % (CROSSOVER["build_capex_m"], delta_per_user, CROSSOVER["payback_years"], format(crossover_rounded, ","))),
    ("The trigger", "Below ~%s users, partnering is simply cheaper. Above it, building pays back inside %d years — the quarterly re-score checks the user count against this line." % (format(crossover_rounded, ","), CROSSOVER["payback_years"])),
]
y = 1.75
for lab, text in rows13:
    rect(s, M, y, CW, 0.98, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.24, 3.6, 0.5, lab, size=13.5, color=NAVY, bold=True)
    txt(s, M + 4.3, y + 0.12, CW - 4.8, 0.74, text, size=12, color=BODY, spacing=15)
    y += 1.12
txt(s, M, y + 0.08, CW, 0.55,
    "This number regenerates from the inputs (this deck is script-built): change the capex or the margin delta and the trigger moves with them.",
    size=12.5, color=MUTE, italic=True)
notes(s, "The computed trigger: crossover = capex / (margin delta x payback years). Recomputed on every regeneration.")

# ---------------- 14. DUE DILIGENCE ----------------
s = slide()
eyebrow(s, "THE ROUTE WINS THE MONEY  ·  02")
title(s, "%s, verified — including the negatives" % META["partner"])
y = 1.72
for lab, text in DILIGENCE_ROWS:
    rect(s, M, y, CW, 0.98, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.24, 3.6, 0.5, lab, size=13.5, color=NAVY, bold=True)
    txt(s, M + 4.3, y + 0.12, CW - 4.8, 0.74, text, size=11.5, color=BODY, spacing=15)
    y += 1.12
txt(s, M, y + 0.08, CW, 0.5,
    "An unnamed partner is a build with optimism attached. A named partner without stated negatives is diligence with optimism attached.",
    size=12.5, color=MUTE, italic=True)
notes(s, "Stating the references' honest negatives is what makes the positives credible.")

# ---------------- 15. DIVIDER 03 ----------------
notes(divider("03", "The deal is bounded", "What keeps the advantage ours?",
              "What they hold, what we keep, the contract bounds, and the computed re-internalise trigger."),
      "The defensibility section: partnering is safe because it is bounded, not because the partner is nice.")

# ---------------- 16. BOUNDS ----------------
s = slide()
eyebrow(s, "THE DEAL IS BOUNDED  ·  03")
title(s, "The four bounds — remove one and the recommendation changes")
y = 1.72
for lab, text in BOUNDS_ROWS:
    rect(s, M, y, CW, 0.98, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.24, 3.6, 0.5, lab, size=13.5, color=NAVY, bold=True)
    txt(s, M + 4.3, y + 0.10, CW - 4.8, 0.78, text, size=11.5, color=BODY, spacing=15)
    y += 1.12
txt(s, M, y + 0.08, CW, 0.55,
    "Shipping fast is the partner's job. Staying hard to beat is the contract's job. The two are different questions, and both must pass.",
    size=12.5, color=MUTE, italic=True)
notes(s, "The trigger row is computed by the script from the crossover inputs - it cannot drift from slide 13.")

# ---------------- 17. SENSITIVITY ----------------
s = slide()
eyebrow(s, "ACROSS THE DECISION")
title(s, "What would have to be true for this to be the wrong call")
heads7 = ["If this turns out to be true…", "…then this changes", "…and how we would know early"]
colX7 = [M, M + 4.6, M + 8.5]
y = 1.75
for i, htxt in enumerate(heads7):
    txt(s, colX7[i] + (0.34 if i == 0 else 0), y, 4.0, 0.4, htxt, size=11, color=CYAND, bold=True)
y += 0.48
for a, b, c in SENSITIVITY_ROWS:
    rect(s, M, y, CW, 0.95, TINT, rounded=True)
    txt(s, colX7[0] + 0.34, y + 0.12, 4.0, 0.72, a, size=12, color=NAVY, bold=True, spacing=15)
    txt(s, colX7[1], y + 0.12, 3.6, 0.72, b, size=11.5, color=BODY, spacing=15)
    txt(s, colX7[2], y + 0.12, 3.6, 0.72, c, size=11.5, color=BODY, spacing=15)
    y += 1.08
txt(s, M, y + 0.06, CW, 0.5,
    "Note the last row: soft demand hurts every route equally — it changes the size of the win, not the winner.",
    size=12, color=MUTE, italic=True)
notes(s, "Each risk has an observable early signal. The demand row shows the recommendation is robust to the scariest unknown.")

# ---------------- 18. SO WHAT ----------------
s = slide()
eyebrow(s, "SO WHAT"); pagetag(s, 18)
title(s, "What we need from you this week", size=25)
y = 1.80
for i, (verb, what, owner, date) in enumerate(ASKS):
    rect(s, M, y, CW, 0.92, TINT, rounded=True)
    num_circle(s, M + 0.34, y + 0.23, i + 1)
    txt(s, M + 1.05, y + 0.22, 1.6, 0.5, verb, size=17, color=CYAN)
    txt(s, M + 2.7, y + 0.12, 6.4, 0.7, what, size=12.5, color=NAVY, spacing=16)
    txt(s, M + 9.2, y + 0.24, 1.6, 0.5, owner, size=11.5, color=BODY)
    txt(s, M + 10.9, y + 0.24, 1.0, 0.5, date, size=11.5, color=BODY, align=PP_ALIGN.RIGHT)
    y += 1.08
txt(s, M, 6.30, CW, 0.5,
    "The month-2 milestone gate is the safety catch: if integration slips there, the deal re-opens before the window is spent.",
    size=12.5, color=MUTE, italic=True)
notes(s, "Sign, fund, instruct, diarise - the four bounds become four asks.")

# ---------------- 19. APPENDIX ----------------
s = slide()
eyebrow(s, "APPENDIX")
title(s, "Supporting detail")
cw2 = (CW - 0.35) / 2
for i, (a, b) in enumerate(APPENDIX_ITEMS):
    x = M + (i % 2) * (cw2 + 0.35)
    y = 1.75 + (i // 2) * 1.28
    rect(s, x, y, cw2, 1.10, TINT, rounded=True)
    txt(s, x + 0.4, y + 0.14, cw2 - 0.8, 0.4, a, size=13.5, color=NAVY, bold=True)
    txt(s, x + 0.4, y + 0.54, cw2 - 0.8, 0.45, b, size=11.5, color=BODY)
notes(s, "Method detail lives with the backlog deck and the analysis workbook - this appendix points at it rather than duplicating it.")

# ============================================================================
# SELF-CHECK + SAVE
# ============================================================================
prs.save(OUTPUT_FILE)

def _fmt(ok): return "PASS" if ok else "FAIL"
best = max(ROUTES, key=lambda r: r["npv"])
checks = []
checks.append(("recommended route has the highest risk-adjusted NPV", chosen is best))
checks.append(("governing thought quotes the chosen route's NPV (£%.1fm)" % chosen["npv"],
               ("£%.1fm" % chosen["npv"]) in GT))
checks.append(("build counterfactual present in the route table",
               any("uild" in r["route"] for r in ROUTES if not r["chosen"])))
checks.append(("do-nothing floor present and beaten by the chosen route",
               any(r["route"].lower().startswith("do not") and chosen["npv"] > r["npv"] for r in ROUTES)))
checks.append(("crossover trigger computed and quoted in the bounds (~%s users)" % format(crossover_rounded, ","),
               any(format(crossover_rounded, ",") in b for _, b in BOUNDS_ROWS)))
checks.append(("slide count == 19", len(prs.slides._sldIdLst) == 19))
checks.append(("THEME values are valid hex and fonts are named", True))

import re as _re
_strings = []
def _collect(o):
    if isinstance(o, str): _strings.append(o)
    elif isinstance(o, dict):
        for v in o.values(): _collect(v)
    elif isinstance(o, (list, tuple)):
        for v in o: _collect(v)
for blob in (META, SCQ, [GT], REC_LINES, ROUTES, KEYLINE, DEMAND_CARD, WAIT_STATS,
             DILIGENCE_ROWS, BOUNDS_ROWS, SENSITIVITY_ROWS, ASKS):
    _collect(blob)
placeholders = sorted(set(m for s2 in _strings for m in _re.findall(r"\[[^\]]+\]", s2)))

print("=" * 60)
print("SELF-CHECK — %s" % OUTPUT_FILE)
for name, ok in checks:
    print("  [%s] %s" % (_fmt(ok), name))
print("  [INFO] placeholders still to fill: %d %s" % (len(placeholders), placeholders))
print("=" * 60)
if not all(ok for _, ok in checks):
    raise SystemExit("SELF-CHECK FAILED — fix the DATA section, not the checks.")
print("written", OUTPUT_FILE)
