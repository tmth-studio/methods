#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PRESENTED DECK GENERATOR — text-light Minto slides for a speaker in the room.

Produces a .pptx where the PRESENTER carries the argument: headline claims,
one exhibit or at most three support points per slide, and speaker notes that
hold the narrative. For a deck that must travel without a presenter, use
generate_standalone_deck.py instead — it takes the identical DATA.

HOW TO USE (in Copilot, Claude, or any assistant with a code interpreter):
  1. Attach this file. Say: "Run this script and give me the .pptx it creates."
     It ships with fictional demonstration content (Meridian Office Services).
  2. To build your own deck: the assistant follows the STAGED PROMPTS below
     and edits ONLY the DATA section. Everything under BUILD — DO NOT EDIT
     stays. Re-run after every edit; the SELF-CHECK refuses to write a file
     while any rule fails.

THE SHARED SCHEMA: the PYRAMID data below is identical to the DATA sections of
generate_pyramid_memo.py and generate_standalone_deck.py. Fill it once, paste
it between the DATA markers of all three, and one argument drives the memo
and both decks. (Each script sets its own output filename in the build
section, so a pasted DATA block never renames the artefact.) Generator-
specific registers sit at the END of each DATA section (the decks:
EXHIBIT_REGISTER + MAX_VINTAGE_MONTHS; the memo: EVIDENCE_REGISTER) — when
pasting the shared block between scripts, keep each script's own registers
in place.

STAGED PROMPTS — the assistant follows these in order and STOPS at each gate.

  STAGE 0 — APPLY THE CLIENT'S BRAND (optional). All brand values live in
  the THEME dict at the top of the DATA section; the shipped default is the
  neutral TMTH palette. Two routes to a client brand: (i) paste the client's
  palette hexes and font names into THEME; (ii) attach any existing client
  .pptx and run extract_theme.py (same repository) — it prints a
  ready-to-paste THEME dict with a suggested role mapping. Client brand
  arrives at runtime; never commit a client THEME to the public repository.

  STAGE 1 — FRAME. Ask, one question at a time: who is the audience; what
  decision is wanted; which of the four patterns (giving direction / seeking
  approval / explaining how / choosing); confirm the author WILL present —
  if the deck travels without a speaker, switch to the standalone generator.
  GATE: get a yes on the frame.

  STAGE 2 — SCQA AND GOVERNING THOUGHT. Draft Situation, Complication, one
  Question, and the governing thought in three drafts (WHAT → + SO-WHAT →
  25 words or fewer, active, verb first). Propose the tone order with a
  reason. GATE: show the drafts, get a yes.

  STAGE 3 — PYRAMID. 2–4 key-line arguments as full-sentence assertions;
  name the ordering principle; MECE test (one exhibit proving two arguments =
  merge; if all true the audience MUST accept the answer, else something is
  missing); 2–5 supports per argument from evidence the user supplies.
  GATE: show the pyramid + audit, get a yes.

  STAGE 4 — EXHIBIT REUSE. Attach the engagement's prior decks. Catalogue
  the existing exhibits — for each: the MESSAGE it proves (a claim, not a
  topic), its chart form, its source deck, and its data vintage (the date
  of the underlying data, not the deck date). For every key-line point
  needing an exhibit, propose REUSE / ADAPT (say exactly what changes) /
  NEW, and fill the EXHIBIT_REGISTER. Staleness guard: any REUSE or ADAPT
  whose data vintage is older than MAX_VINTAGE_MONTHS (default 6) FAILS the
  build unless flagged vintage_accepted with a stated reason. Only NEW and
  ADAPT exhibits get design effort in Stage 5. GATE: show the register,
  get a yes.

  STAGE 5 — EXHIBITS AND NOTES. For each support: write the claim first, name
  the comparison it implies (component / item / time series / frequency /
  correlation), pick the chart kind from the comparison, and supply the real
  data — never invented. Give each exhibit its register exhibit_id. Keep
  bullets to three, under 20 words each. Write speaker notes that carry the
  narrative the slide deliberately leaves off. GATE: show the exhibit spec
  table, get a yes.

  STAGE 6 — FILL AND RUN. Write the approved content into DATA. Replace every
  square-bracketed placeholder — the build fails while any remain. Run the
  script; hand back the .pptx and the self-check report.

ENGAGEMENT ASSETS: a filled EXHIBIT_REGISTER is client work product. It
lives in the client tenant. Never commit a filled register to the public
repository — the shipped register holds fictional demonstration content only.

VERSION NOTES
  v1.1 — 20 Aug 2026: brand moved from code to data (THEME dict, neutral
         TMTH default; validation gate on hex/fonts; dark-slide tints now
         derived from the dominant colour). Exhibit-reuse stage added with
         the EXHIBIT_REGISTER, exhibit_id linking, and a data-vintage
         staleness guard (MAX_VINTAGE_MONTHS, default 6).
  v1.0 — 20 Aug 2026: first published version.

Requires: python-pptx.  Output: presented-deck.pptx
"""

# ============================================================================
# DATA — EDIT THIS SECTION ONLY  (shared PYRAMID schema — see docstring)
# All demonstration content is fictional (Meridian Office Services).
# ============================================================================

# --- THEME: all brand values live here, as data. Ship = neutral TMTH palette.
# To apply a client's brand: (i) paste their palette hexes + font names below,
# or (ii) attach any existing client .pptx and run extract_theme.py (same
# repository) — it prints a ready-to-paste THEME dict. Hex: 6 chars, no '#',
# no alpha. Never commit a client THEME to the public repository.
THEME = {
    "dominant":      "2F4858",  # lead brand colour — cover, dividers, the answer, table headers
    "accent":        "4E7C90",  # secondary brand colour — slide titles, line charts
    "highlight":     "C9A227",  # the one-highlight colour — eyebrows, numbers, highlighted bars
    "card_fill":     "EEF3F6",  # card / panel background tint
    "ink":           "1F2A31",  # near-black headings ink
    "ink_soft":      "3C4C57",  # body text ink
    "muted":         "6B7A85",  # secondary text, sources, footers
    "chart_compare": "B9C6CD",  # non-highlighted chart bars / comparison elements
    "font_heading":  "Cambria",
    "font_body":     "Calibri",
}

META = {
    "title": "Win renewals on response time, not price",
    "subtitle": "Recommendation — SME contract churn",
    "author_line": "Strategy team",
    "date": "20 Aug 2026",
    "audience": "Managing director and the executive committee",
    "presence": "live",              # this generator assumes a speaker
    "archetype": "recommendation",   # used by the memo generator; ignored here
    "footer_tag": "Demonstration content — every name and number is fictional",
}

SCQA = {
    "tone_order": "SCA",  # SCA standard | ASC direct | CSA concerned
    "situation": ("Meridian serves 1,200 SME service contracts across the "
                  "South East, worth £5.1M of annual revenue. Renewal pricing "
                  "has been flat for two years."),
    "complication": ("Churn has risen from 22% to 35% in three years — a "
                     "£1.8M annual revenue impact. The lost-deal evidence "
                     "points away from price: objections fell 30% while churn "
                     "climbed."),
    "question": ("How does Meridian win renewals in a market that now judges "
                 "providers on response time?"),
}

GOVERNING_THOUGHT = ("Reposition Meridian around a response-time guarantee — "
                     "the factor deciding renewals — recovering churn to 25% "
                     "and £1.2M annual revenue within four quarters.")

KEY_LINE_ORDER = "degree"  # "time" | "structure" | "degree"

KEY_LINE = [
    {
        "tag": "DIAGNOSIS",
        "assertion": "Churn is perception-driven, not price-driven",
        "why_it_matters": ("Every remedy on the table prices the wrong "
                           "problem. The evidence separates what customers "
                           "say from what they do — and it clears price as "
                           "the cause."),
        "supports": [
            {
                "assertion": ("Current customers renew at stable rates while "
                              "competitor-bid churn climbs"),
                "prose": ("Existing-contract churn has held between 21% and "
                          "23% for eight quarters. A price problem hits "
                          "paying customers first, and it has not. The rise "
                          "is concentrated in renewals where a competitor "
                          "bid is on the table."),
                "bullets": [
                    "Existing-contract churn steady at 21–23% for eight quarters",
                    "Price pain reaches current payers first — it has not",
                    "The rise sits in renewals facing a competitor bid",
                ],
                "source": "Contract system renewal records, Q3 2024 to Q2 2026, n=1,180",
                "exhibit": None,
                "notes": ("Land the distinction between the two churn "
                          "populations before anyone mentions price. The "
                          "whole argument turns on it."),
            },
            {
                "assertion": ("Price objections fell 30% while churn rose 13 "
                              "points"),
                "prose": ("Recorded price objections fell from 46 to 32 a "
                          "quarter across three years. If price drove the "
                          "losses, objections would rise with churn. They "
                          "moved the opposite way."),
                "bullets": [
                    "Objections fell from 46 to 32 a quarter",
                    "Churn rose 13 points across the same period",
                    "The two series move in opposite directions",
                ],
                "source": "CRM objection log, 2023 to 2026",
                "exhibit": {"exhibit_id": "EX-OBJECTIONS",
                            "kind": "column",
                            "unit": "Price objections per quarter",
                            "categories": ["2023", "2024", "2025", "2026"],
                            "values": [46, 41, 37, 32],
                            "highlight": 3,
                            "source": "CRM objection log, 2023 to 2026"},
                "notes": ("If price were the cause, this line would rise "
                          "with churn. Let the divergence do the work — no "
                          "editorialising needed."),
            },
            {
                "assertion": ("All 10 lost Q4 deals cite the competitor "
                              "guarantee and none cite price"),
                "prose": ("Every lost-deal review from the last quarter "
                          "names the competitor's four-hour response "
                          "guarantee. Price appears in none of them. The "
                          "market has changed what it judges providers on."),
                "bullets": [
                    "10 of 10 reviews cite the response guarantee",
                    "Zero reviews cite price",
                    "Scope and relationship reasons trail far behind",
                ],
                "source": "Lost-deal reviews, Oct to Dec 2025, n=10",
                "exhibit": {"exhibit_id": "EX-LOSTDEAL",
                            "kind": "bar",
                            "unit": "Mentions across 10 lost-deal reviews",
                            "categories": ["Competitor response guarantee",
                                           "Service scope",
                                           "Relationship change",
                                           "Price"],
                            "values": [10, 3, 2, 0],
                            "highlight": 0,
                            "source": "Lost-deal reviews, Oct to Dec 2025, n=10"},
                "notes": ("This is the exhibit the sceptics remember. Ten "
                          "out of ten is unusual — say so, and say the n is "
                          "small but the signal is unanimous."),
            },
        ],
        "transition": ("If price is not the lever, the lost-deal records say "
                       "what is — and it is buildable."),
    },
    {
        "tag": "REMEDY",
        "assertion": ("A four-hour response guarantee wins the frame the "
                      "market now uses"),
        "why_it_matters": ("The competitor changed the question renewals are "
                           "judged on. Matching the guarantee answers that "
                           "question directly — and Meridian starts closer "
                           "to it than the sales narrative admits."),
        "supports": [
            {
                "assertion": ("The competitor guarantee reset what buyers ask "
                              "at renewal"),
                "prose": ("Renewal conversations now open with response "
                          "commitments, not price. Account managers report "
                          "the guarantee is raised by the customer in most "
                          "competitive renewals. The frame has moved and the "
                          "proposition has not moved with it."),
                "bullets": [
                    "Renewal conversations now open with response commitments",
                    "Customers raise the guarantee unprompted",
                    "The proposition has not moved with the frame",
                ],
                "source": "Account team renewal debriefs, Q1 to Q2 2026",
                "exhibit": None,
                "notes": ("Anchor this in the account managers' own words — "
                          "the evidence is qualitative and honest about it."),
            },
            {
                "assertion": ("Meridian already meets four hours on 78% of "
                              "callouts"),
                "prose": ("Dispatch records show 78% of 2025 callouts were "
                          "answered inside four hours without any guarantee "
                          "in place. The gap to close is the remaining 22% — "
                          "evenings, weekends and multi-site jobs. That gap "
                          "is known and priceable."),
                "bullets": [
                    "78% of 2025 callouts already inside four hours",
                    "The gap is evenings, weekends and multi-site jobs",
                    "The gap is known and priceable",
                ],
                "source": "Dispatch system records, Jan to Dec 2025, n=8,400",
                "exhibit": {"exhibit_id": "EX-CALLOUTS",
                            "kind": "bar",
                            "unit": "% of 2025 callouts",
                            "categories": ["Met within four hours today",
                                           "Requires new cover"],
                            "values": [78, 22],
                            "highlight": 0,
                            "source": "Dispatch system records, 2025, n=8,400"},
                "notes": ("The surprise is how close we already are. The "
                          "guarantee prices the last 22%, not the whole "
                          "operation."),
            },
            {
                "assertion": ("The guarantee costs £350k a year to stand up"),
                "prose": ("Three cost lines close the 22% gap: extended "
                          "engineering cover, a triage desk, and a service-"
                          "credit reserve priced at a 2% breach rate. The "
                          "total is £350k a year, fixed."),
                "bullets": [
                    "Extended engineering cover — £190k",
                    "Triage desk and dispatch tooling — £95k",
                    "Service-credit reserve at 2% breach — £65k",
                ],
                "source": "Operations costing model v2, Jul 2026",
                "exhibit": {"exhibit_id": "EX-COSTS",
                            "kind": "table",
                            "headers": ["Cost line", "Annual cost"],
                            "rows": [
                                ["Evening and weekend engineering cover", "£190k"],
                                ["Triage desk and dispatch tooling", "£95k"],
                                ["Service-credit reserve at 2% breach rate", "£65k"],
                                ["Total", "£350k"],
                            ],
                            "source": "Operations costing model v2, Jul 2026"},
                "notes": ("The reserve line is the honest one — a guarantee "
                          "without a priced breach rate is a slogan."),
            },
        ],
        "transition": ("A guarantee Meridian can meet at a known cost must "
                       "still pay. It does — on pessimistic assumptions."),
    },
    {
        "tag": "ECONOMICS",
        "assertion": ("The programme pays back within four quarters on the "
                      "central case"),
        "why_it_matters": ("The case is priced in ranges, not a point value. "
                           "Even the pessimistic case clears the outlay "
                           "inside five quarters, and the staged spend caps "
                           "the downside."),
        "supports": [
            {
                "assertion": ("Recovered renewals return £0.9M to £1.6M a "
                              "year against the £350k outlay"),
                "prose": ("The scenario table runs churn recovery at three "
                          "levels. The central case recovers churn to 27% "
                          "and returns £1.2M a year — payback inside four "
                          "quarters. The pessimistic case still pays back "
                          "in five."),
                "bullets": [
                    "Central case: churn to 27%, £1.2M a year recovered",
                    "Pessimistic case still pays back in five quarters",
                    "All three cases clear the £350k outlay",
                ],
                "source": "Renewal economics model, scenario run Jul 2026",
                "exhibit": {"exhibit_id": "EX-SCENARIOS",
                            "kind": "table",
                            "headers": ["Scenario", "Churn recovered to",
                                        "Revenue recovered", "Payback"],
                            "rows": [
                                ["Pessimistic", "31%", "£0.9M a year", "5 quarters"],
                                ["Central", "27%", "£1.2M a year", "4 quarters"],
                                ["Optimistic", "25%", "£1.6M a year", "3 quarters"],
                            ],
                            "source": "Renewal economics model, Jul 2026"},
                "notes": ("Never present the central case alone. The range "
                          "is the credibility."),
            },
            {
                "assertion": ("The 90-day gate caps the downside at £120k"),
                "prose": ("Spend is staged: £120k reaches the field before "
                          "the 90-day review. If lost-deal citations of the "
                          "competitor guarantee have not fallen below five a "
                          "quarter by then, the programme stops before the "
                          "remaining spend commits."),
                "bullets": [
                    "Only £120k committed before the 90-day review",
                    "Stop trigger: citations still above five a quarter",
                    "The remaining spend never commits on a failed signal",
                ],
                "source": "Programme spend schedule, draft Jul 2026",
                "exhibit": None,
                "notes": ("The gate is what makes this approvable — the "
                          "decision is reversible at a known cost."),
            },
        ],
        "transition": None,
    },
]

CONCLUSION = {
    "restate": ("The renewal problem is a framing problem, and the guarantee "
                "answers it — churn back to 25% and £1.2M of revenue "
                "recovered within four quarters."),
    "key_assumption": ("This rests on the market's response-time frame "
                       "persisting. If lost-deal citations shift away from "
                       "it, we review at the 90-day gate."),
    "next_action": ("The operations director stands up the build team on "
                    "1 Oct 2026."),
}

ASK = {
    "decision": ("Approve the response-time guarantee programme at £350k a "
                 "year"),
    "gate": ("Stop at 90 days if lost-deal citations of the competitor "
             "guarantee have not fallen below five a quarter"),
    "owner": "Operations director",
    "date": "1 Oct 2026",
}

NEXT_STEPS = [
    ("Confirm the service-credit terms with legal", "Commercial lead", "12 Sep 2026"),
    ("Brief the account teams on the renewal script", "Sales director", "19 Sep 2026"),
]

# --- EXHIBIT REUSE (Stage 4) — one row per exhibit this deck uses. Every
# support exhibit's exhibit_id must appear here. decision: "REUSE" (as-is
# from a prior deck) | "ADAPT" ('changes' required) | "NEW". data_vintage =
# date of the UNDERLYING data (YYYY-MM), not the deck date. Staleness guard:
# REUSE/ADAPT rows older than MAX_VINTAGE_MONTHS fail the build unless
# vintage_accepted is True with a stated vintage_reason. NEW rows are exempt.
# ENGAGEMENT ASSET: never commit a filled register to the public repository.
MAX_VINTAGE_MONTHS = 6

EXHIBIT_REGISTER = [
    {"exhibit_id": "EX-OBJECTIONS",
     "message": ("Price objections fell while churn rose — the two series "
                 "diverge"),
     "chart_form": "column", "source_deck": "Q2 renewal review deck",
     "data_vintage": "2026-07", "decision": "ADAPT",
     "changes": ("Extend the series to 2026 and move the highlight to the "
                 "latest year"),
     "vintage_accepted": False, "vintage_reason": ""},
    {"exhibit_id": "EX-LOSTDEAL",
     "message": ("All 10 lost Q4 deals cite the competitor guarantee; none "
                 "cite price"),
     "chart_form": "bar", "source_deck": "",
     "data_vintage": "2025-12", "decision": "NEW", "changes": "",
     "vintage_accepted": False, "vintage_reason": ""},
    {"exhibit_id": "EX-CALLOUTS",
     "message": "Meridian already meets four hours on 78% of callouts",
     "chart_form": "bar", "source_deck": "",
     "data_vintage": "2026-02", "decision": "NEW", "changes": "",
     "vintage_accepted": False, "vintage_reason": ""},
    {"exhibit_id": "EX-COSTS",
     "message": "The guarantee costs £350k a year, in three fixed lines",
     "chart_form": "table", "source_deck": "Operations costing pack v2",
     "data_vintage": "2026-01", "decision": "REUSE", "changes": "",
     "vintage_accepted": True,
     "vintage_reason": ("Costing re-validated Jul 2026 against supplier "
                        "quotes; underlying rates unchanged")},
    {"exhibit_id": "EX-SCENARIOS",
     "message": ("Recovered renewals return £0.9M–£1.6M a year against the "
                 "£350k outlay"),
     "chart_form": "table", "source_deck": "Renewal economics model pack",
     "data_vintage": "2026-07", "decision": "REUSE", "changes": "",
     "vintage_accepted": False, "vintage_reason": ""},
]

# ============================================================================
# BUILD — DO NOT EDIT BELOW THIS LINE
# ============================================================================
import datetime
import re
import sys

OUTPUT_FILE = "presented-deck.pptx"

# ---------------------------------------------------------------- self-check
BANNED_LABELS = {"background", "analysis", "overview", "introduction",
                 "context", "summary", "next steps", "recommendation",
                 "findings", "update", "situation", "options", "approach",
                 "details", "conclusion"}

def _assertion_defect(text):
    t = (text or "").strip()
    words = t.split()
    if len(words) < 5:
        return "under 5 words — reads as a label, not a claim"
    if t.rstrip(".").lower() in BANNED_LABELS:
        return "banned label heading"
    if t.endswith(":"):
        return "ends in a colon — a label, not a claim"
    tail = [w for w in words[1:] if w[:1].isalpha()]
    caps = sum(1 for w in tail if w[:1].isupper())
    if tail and caps > 0.6 * len(tail):
        return "Title Case — use sentence case"
    return None

def _collect_strings(obj, out):
    if isinstance(obj, str):
        out.append(obj)
    elif isinstance(obj, dict):
        for v in obj.values():
            _collect_strings(v, out)
    elif isinstance(obj, (list, tuple)):
        for v in obj:
            _collect_strings(v, out)

def _valid_exhibit(ex):
    if ex is None:
        return None
    if not isinstance(ex, dict):
        return "exhibit is not a dict"
    kind = ex.get("kind")
    if kind in ("bar", "column", "line"):
        cats, vals = ex.get("categories"), ex.get("values")
        if not cats or not vals or len(cats) != len(vals):
            return "chart categories/values missing or mismatched"
        hl = ex.get("highlight")
        if hl is not None and not (0 <= hl < len(vals)):
            return "highlight index out of range"
    elif kind == "table":
        if not ex.get("headers") or not ex.get("rows"):
            return "table needs headers and rows"
    else:
        return "exhibit kind must be bar / column / line / table"
    if not ex.get("source"):
        return "exhibit has no source"
    return None

EXHIBIT_DECISIONS = ("REUSE", "ADAPT", "NEW")

def _theme_defects():
    defects = []
    for k, v in THEME.items():
        if k.startswith("font_"):
            if not str(v).strip():
                defects.append("THEME['%s'] is empty — name a real font" % k)
        elif not re.fullmatch(r"[0-9A-Fa-f]{6}", str(v)):
            defects.append("THEME['%s'] = %r is not a 6-char hex "
                           "(no '#', no alpha)" % (k, v))
    return defects

def _vintage_age_months(vintage):
    """Age of a YYYY-MM vintage in months, or None if malformed."""
    m = re.fullmatch(r"(\d{4})-(\d{2})", str(vintage).strip())
    if not m:
        return None
    y, mo = int(m.group(1)), int(m.group(2))
    if not 1 <= mo <= 12:
        return None
    today = datetime.date.today()
    return (today.year - y) * 12 + (today.month - mo)

def _exhibit_register_defects():
    """Structural defects and the staleness guard, separately named."""
    struct, stale = [], []
    seen = set()
    for r in EXHIBIT_REGISTER:
        rid = str(r.get("exhibit_id", "")).strip() or "?"
        if rid == "?":
            struct.append("a register row has no exhibit_id")
        elif rid in seen:
            struct.append("duplicate exhibit_id %s" % rid)
        seen.add(rid)
        if not str(r.get("message", "")).strip():
            struct.append("%s has no message — name the claim it proves"
                          % rid)
        dec = r.get("decision")
        if dec not in EXHIBIT_DECISIONS:
            struct.append("%s decision must be REUSE / ADAPT / NEW" % rid)
        if dec == "ADAPT" and not str(r.get("changes", "")).strip():
            struct.append("%s is ADAPT with no 'changes' — say what "
                          "changes" % rid)
        if r.get("vintage_accepted") and not str(
                r.get("vintage_reason", "")).strip():
            struct.append("%s is vintage_accepted with no stated reason"
                          % rid)
        age = _vintage_age_months(r.get("data_vintage", ""))
        if dec in ("REUSE", "ADAPT"):
            if not str(r.get("source_deck", "")).strip():
                struct.append("%s is %s with no source_deck" % (rid, dec))
            if age is None:
                struct.append("%s data_vintage must be YYYY-MM (date of "
                              "the underlying data)" % rid)
            elif age > MAX_VINTAGE_MONTHS and not (
                    r.get("vintage_accepted")
                    and str(r.get("vintage_reason", "")).strip()):
                stale.append("%s: data vintage %s is %d months old (max "
                             "%d) — rebuild as NEW or set vintage_accepted "
                             "with a reason"
                             % (rid, r["data_vintage"], age,
                                MAX_VINTAGE_MONTHS))
    return struct, stale

def run_self_check():
    checks, warnings = [], []
    td = _theme_defects()
    checks.append(("THEME: colours are 6-char hex and fonts are named"
                   + ("" if not td else " — " + "; ".join(td)), not td))
    gt_words = len(GOVERNING_THOUGHT.split())
    checks.append(("governing thought is 25 words or fewer (%d)" % gt_words,
                   0 < gt_words <= 25))
    checks.append(("SCQA complete: situation, complication, question, answer",
                   all([SCQA.get("situation"), SCQA.get("complication"),
                        SCQA.get("question"), GOVERNING_THOUGHT])))
    low_sit = " " + SCQA["situation"].lower()
    checks.append(("situation carries no 'but' / 'however'",
                   " but " not in low_sit and "however" not in low_sit))
    checks.append(("exactly one question, ending in a question mark",
                   SCQA["question"].count("?") == 1
                   and SCQA["question"].strip().endswith("?")))
    checks.append(("key line holds 2–4 arguments (%d)" % len(KEY_LINE),
                   2 <= len(KEY_LINE) <= 4))
    checks.append(("ordering principle named: time / structure / degree",
                   KEY_LINE_ORDER in ("time", "structure", "degree")))

    label_defects = []
    for kl in KEY_LINE:
        d = _assertion_defect(kl.get("assertion", ""))
        if d:
            label_defects.append("key line '%s…': %s"
                                 % (kl.get("assertion", "")[:40], d))
        for sp in kl.get("supports", []):
            d = _assertion_defect(sp.get("assertion", ""))
            if d:
                label_defects.append("support '%s…': %s"
                                     % (sp.get("assertion", "")[:40], d))
    checks.append(("every slide title is a full-sentence claim, not a label"
                   + ("" if not label_defects else
                      " — " + "; ".join(label_defects)),
                   not label_defects))

    checks.append(("every key-line argument carries 2–5 supports",
                   all(2 <= len(kl.get("supports", [])) <= 5
                       for kl in KEY_LINE)))
    kl_assertions = [kl["assertion"].strip().lower() for kl in KEY_LINE]
    checks.append(("no two key-line arguments make the same claim (MECE "
                   "overlap)", len(set(kl_assertions)) == len(kl_assertions)))
    checks.append(("every support cites a source",
                   all(sp.get("source") for kl in KEY_LINE
                       for sp in kl["supports"])))

    # presented-deck discipline: text-light bodies, notes carry the narrative
    body_defects = []
    for kl in KEY_LINE:
        for sp in kl["supports"]:
            b = sp.get("bullets") or []
            if not sp.get("exhibit") and not b:
                body_defects.append("'%s…': no exhibit and no bullets"
                                    % sp["assertion"][:35])
            if len(b) > 3:
                body_defects.append("'%s…': more than 3 bullets"
                                    % sp["assertion"][:35])
            for x in b:
                if len(x.split()) > 20:
                    body_defects.append("'%s…': bullet over 20 words"
                                        % sp["assertion"][:35])
    checks.append(("text-light bodies: at most 3 bullets, 20 words each, or "
                   "an exhibit" + ("" if not body_defects else
                                   " — " + "; ".join(body_defects)),
                   not body_defects))
    checks.append(("every support carries speaker notes (the presenter holds "
                   "the narrative)",
                   all((sp.get("notes") or "").strip() for kl in KEY_LINE
                       for sp in kl["supports"])))

    ex_defects = []
    for kl in KEY_LINE:
        for sp in kl["supports"]:
            d = _valid_exhibit(sp.get("exhibit"))
            if d:
                ex_defects.append("'%s…': %s" % (sp["assertion"][:35], d))
    checks.append(("every exhibit is well-formed with a source"
                   + ("" if not ex_defects else " — " + "; ".join(ex_defects)),
                   not ex_defects))

    # --- exhibit reuse register (Stage 4)
    reg_struct, reg_stale = _exhibit_register_defects()
    checks.append(("exhibit register rows are well-formed: message, valid "
                   "decision, ADAPT changes, REUSE/ADAPT source + vintage"
                   + ("" if not reg_struct else
                      " — " + "; ".join(reg_struct)), not reg_struct))
    checks.append(("staleness guard: no REUSE/ADAPT older than %d months "
                   "without an accepted vintage" % MAX_VINTAGE_MONTHS
                   + ("" if not reg_stale else
                      " — " + "; ".join(reg_stale)), not reg_stale))
    reg_ids = set(str(r.get("exhibit_id", "")).strip()
                  for r in EXHIBIT_REGISTER)
    link_defects = []
    for kl in KEY_LINE:
        for sp in kl["supports"]:
            ex = sp.get("exhibit")
            if ex is None:
                continue
            eid = str(ex.get("exhibit_id", "")).strip()
            if not eid:
                link_defects.append("'%s…' exhibit has no exhibit_id"
                                    % sp["assertion"][:35])
            elif eid not in reg_ids:
                link_defects.append("'%s…' exhibit_id %s has no register "
                                    "row" % (sp["assertion"][:35], eid))
    checks.append(("every support exhibit carries an exhibit_id with a "
                   "register row (REUSE / ADAPT / NEW decided)"
                   + ("" if not link_defects else
                      " — " + "; ".join(link_defects)), not link_defects))

    checks.append(("single ask with decision, gate, owner and date",
                   all([ASK.get("decision"), ASK.get("gate"),
                        ASK.get("owner"), ASK.get("date")])))

    strings = []
    for blob in (META, SCQA, GOVERNING_THOUGHT, KEY_LINE, CONCLUSION, ASK,
                 NEXT_STEPS):
        _collect_strings(blob, strings)
    placeholders = sorted(set(m for s in strings
                              for m in re.findall(r"\[[^\]]+\]", s)))
    checks.append(("no bracketed placeholders remain"
                   + ("" if not placeholders else
                      " — replace: " + ", ".join(placeholders)),
                   not placeholders))
    return checks, warnings

def print_report(checks, warnings, extra=(), wrote=False):
    print("=" * 64)
    print("SELF-CHECK — %s" % OUTPUT_FILE)
    for name, ok in list(checks) + list(extra):
        print("  [%s] %s" % ("PASS" if ok else "FAIL", name))
    for w in warnings:
        print("  [WARN] %s" % w)
    print("=" * 64)
    if not all(ok for _, ok in list(checks) + list(extra)):
        raise SystemExit("SELF-CHECK FAILED — fix the DATA section, not the "
                         "checks.")
    if wrote:
        print("written", OUTPUT_FILE)

_checks, _warnings = run_self_check()
if not all(ok for _, ok in _checks):
    print_report(_checks, _warnings)
    sys.exit(1)

# ---------------------------------------------------------------- pptx build
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

def _rgb(h):
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def _blend(base_hex, mix_hex, f):
    """f = share of mix_hex blended into base_hex."""
    b = [int(base_hex[i:i + 2], 16) for i in (0, 2, 4)]
    m = [int(mix_hex[i:i + 2], 16) for i in (0, 2, 4)]
    return RGBColor(*(round(b[i] * (1 - f) + m[i] * f) for i in range(3)))

SLATE = _rgb(THEME["dominant"])      # dominant — cover, dividers, answer
BLUE = _rgb(THEME["accent"])         # headers, chart series
GOLD = _rgb(THEME["highlight"])      # accent — numbers, eyebrows
TINT = _rgb(THEME["card_fill"])      # card fill
INK = _rgb(THEME["ink"])             # body headings
BODY = _rgb(THEME["ink_soft"])
MUTE = _rgb(THEME["muted"])
GREYC = _rgb(THEME["chart_compare"])  # comparison chart elements
WHITE = RGBColor(0xFF, 0xFF, 0xFF)   # functional — not brand
RULE = _blend("FFFFFF", THEME["dominant"], 0.18)
# derived tints of the dominant colour, for text on dark slides
DOM_LT = _blend(THEME["dominant"], "FFFFFF", 0.75)   # light body on dark
DOM_MID = _blend(THEME["dominant"], "FFFFFF", 0.50)  # metadata on dark
DOM_NUM = _blend(THEME["dominant"], "FFFFFF", 0.22)  # watermark numerals
HEAD_FONT, BODY_FONT = THEME["font_heading"], THEME["font_body"]

W, H, M = 13.333, 7.5, 0.62
CW = W - 2 * M

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill, rounded=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    # invisible border: line colour must equal the fill (width 0 still strokes)
    shp.line.color.rgb = fill
    shp.line.width = Pt(0)
    if rounded:
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp

def flood(s, fill):
    """Full-bleed background on dark slides — avoids the 1px edge seam."""
    return rect(s, -0.02, -0.02, W + 0.04, H + 0.04, fill)

def txt(s, x, y, w, h, runs, size=13, color=BODY, bold=False, italic=False,
        align=PP_ALIGN.LEFT, spacing=None, font=BODY_FONT, anchor=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, {})]
    first = True
    for text, opt in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opt.get("align", align)
        if spacing:
            p.line_spacing = Pt(spacing)
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = opt.get("font", font)
        f.size = Pt(opt.get("size", size))
        f.bold = opt.get("bold", bold)
        f.italic = opt.get("italic", italic)
        f.color.rgb = opt.get("color", color)
    return tb

def bullets(s, items, x, y, w, size=14, color=BODY, gap=8):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = "•  " + it
        f = r.font
        f.name = BODY_FONT
        f.size = Pt(size)
        f.color.rgb = color
    return tb

def eyebrow(s, t, y=0.32):
    txt(s, M, y, 10, 0.3, t.upper(), size=11, color=GOLD, bold=True)

def title(s, t, size=24, color=BLUE):
    # slide titles anchor to the TOP of their box so short titles do not float
    txt(s, M, 0.66, CW, 1.05, t, size=size, color=color, font=HEAD_FONT,
        spacing=size + 5, anchor=MSO_ANCHOR.TOP)

def foot(s, t):
    txt(s, M, H - 0.55, CW, 0.3, t, size=9, color=MUTE)

def notes(s, t):
    s.notes_slide.notes_text_frame.text = t

def num_circle(s, x, y, n, d=0.46, fill=GOLD):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d),
                           Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = fill
    c.line.width = Pt(0)
    tf = c.text_frame
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r = p0.add_run()
    r.text = str(n)
    r.font.name = HEAD_FONT
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = WHITE

def add_chart(s, ex, x, y, w, h):
    kinds = {"bar": XL_CHART_TYPE.BAR_CLUSTERED,
             "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
             "line": XL_CHART_TYPE.LINE}
    cd = CategoryChartData()
    cd.categories = [str(c) for c in ex["categories"]]
    cd.add_series(ex.get("unit", "value"), tuple(ex["values"]))
    gf = s.shapes.add_chart(kinds[ex["kind"]], Inches(x), Inches(y),
                            Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.has_legend = False
    ch.has_title = False
    plot = ch.plots[0]
    if ex["kind"] != "line":
        plot.gap_width = 60
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(12)
    dl.font.name = BODY_FONT
    dl.font.bold = True
    dl.font.color.rgb = INK
    if ex["kind"] != "line":
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        ser = plot.series[0]
        hl = ex.get("highlight")
        for i in range(len(ex["values"])):
            pt = ser.points[i]
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = (SLATE if hl is not None
                                             and i == hl else GREYC)
    else:
        ser = plot.series[0]
        ser.format.line.color.rgb = BLUE
        ser.format.line.width = Pt(2.5)
    va = ch.value_axis
    va.visible = False
    va.has_major_gridlines = False
    ca = ch.category_axis
    ca.format.line.fill.background()
    ca.tick_labels.font.size = Pt(11)
    ca.tick_labels.font.name = BODY_FONT
    return gf

def add_table(s, ex, x, y, w):
    headers, rows = ex["headers"], ex["rows"]
    nrows, ncols = 1 + len(rows), len(headers)
    row_h = 0.5
    shp = s.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w),
                             Inches(row_h * nrows))
    tbl = shp.table
    for j, hd in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLATE
        cell.text = str(hd)
        pr = cell.text_frame.paragraphs[0].runs[0].font
        pr.size = Pt(12)
        pr.bold = True
        pr.name = HEAD_FONT
        pr.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            cell = tbl.cell(1 + i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 else TINT
            cell.text = str(v)
            pr = cell.text_frame.paragraphs[0].runs[0].font
            pr.size = Pt(12)
            pr.name = BODY_FONT
            pr.color.rgb = INK
            pr.bold = (i == len(rows) - 1
                       and str(row[0]).lower().startswith("total"))
    return shp

# ---------------- 1. COVER ----------------
s = slide()
flood(s, SLATE)
rect(s, 0, 0, W, 0.16, GOLD)
txt(s, M, 1.5, 10.5, 0.4, META["subtitle"].upper(), size=13, color=GOLD,
    bold=True)
txt(s, M, 2.05, 11.6, 2.2, META["title"], size=44, color=WHITE,
    font=HEAD_FONT, spacing=50)
txt(s, M, 4.35, 11.2, 0.9, GOVERNING_THOUGHT, size=16,
    color=DOM_LT, spacing=22)
txt(s, M, 6.5, 10, 0.4, "%s   ·   %s   ·   %s" % (
    META["author_line"], META["date"], META["audience"]), size=11.5,
    color=DOM_MID)
notes(s, "Presented deck — the speaker carries the argument. Generated by "
         "generate_presentation_deck.py; edit the DATA section and re-run.")

# ---------------- 2. SCQA CONTEXT ----------------
s = slide()
eyebrow(s, "Why we are here")
title(s, "The question this meeting answers")
items = [("SITUATION", SCQA["situation"], 1.30, TINT, BODY),
         ("COMPLICATION", SCQA["complication"], 1.45, TINT, BODY),
         ("QUESTION", SCQA["question"], 1.05, SLATE, WHITE)]
y = 1.85
for lab, text, hh, fill, tcol in items:
    rect(s, M, y, CW, hh, fill, rounded=True)
    txt(s, M + 0.42, y + 0.16, 3.0, 0.3, lab, size=11,
        color=GOLD if fill is SLATE else BLUE, bold=True)
    txt(s, M + 0.42, y + 0.50, CW - 1.0, hh - 0.6, text,
        size=15 if fill is SLATE else 13.5,
        color=tcol, bold=(fill is SLATE), spacing=19)
    y += hh + 0.18
notes(s, "Speak the SCQA — the slide is the reminder, the voice is the "
         "story. Situation and complication are nod-test material only.")

# ---------------- 3. GOVERNING THOUGHT ----------------
s = slide()
flood(s, SLATE)
rect(s, 0, H - 0.16, W, 0.16, GOLD)
txt(s, M, 1.1, 10, 0.4, "THE ANSWER", size=13, color=GOLD, bold=True)
txt(s, M, 2.3, CW, 2.6, GOVERNING_THOUGHT, size=32, color=WHITE,
    font=HEAD_FONT, spacing=42)
notes(s, "Answer first — Minto's rule. Pause after reading it. Everything "
         "after this slide is proof, not suspense.")

# ---------------- 4. KEY LINE OVERVIEW ----------------
s = slide()
eyebrow(s, "The argument in one page")
title(s, "%d reasons, in %s order" % (len(KEY_LINE), KEY_LINE_ORDER))
cwn = (CW - 0.3 * (len(KEY_LINE) - 1)) / len(KEY_LINE)
for i, kl in enumerate(KEY_LINE):
    x = M + i * (cwn + 0.3)
    rect(s, x, 1.85, cwn, 4.1, TINT, rounded=True)
    num_circle(s, x + 0.35, 2.15, i + 1, 0.5)
    txt(s, x + 1.0, 2.22, cwn - 1.3, 0.4, kl["tag"].upper(), size=11,
        color=GOLD, bold=True)
    txt(s, x + 0.35, 2.9, cwn - 0.7, 1.3, kl["assertion"], size=16,
        color=INK, bold=True, font=HEAD_FONT, spacing=21)
    txt(s, x + 0.35, 4.25, cwn - 0.7, 1.5, kl["why_it_matters"], size=11,
        color=BODY, spacing=15)
txt(s, M, 6.25, CW, 0.5, "Each reason stands alone — attack one and the "
    "other %d hold." % (len(KEY_LINE) - 1), size=12, color=MUTE, italic=True)
notes(s, "The key line. Name the ordering principle out loud (%s order) — "
         "it tells the room the list is complete, not brainstormed."
      % KEY_LINE_ORDER)

# ---------------- SECTIONS ----------------
for i, kl in enumerate(KEY_LINE):
    s = slide()
    flood(s, SLATE)
    txt(s, M, 2.0, 2.4, 1.6, "%02d" % (i + 1), size=80,
        color=DOM_NUM, bold=True, font=HEAD_FONT)
    txt(s, M + 2.45, 2.18, 9.6, 0.4, kl["tag"].upper(), size=12, color=GOLD,
        bold=True)
    txt(s, M + 2.45, 2.6, 9.6, 1.4, kl["assertion"], size=30, color=WHITE,
        font=HEAD_FONT, spacing=37)
    txt(s, M + 2.45, 4.05, 9.4, 0.9, kl["why_it_matters"], size=13.5,
        color=DOM_LT, spacing=19)
    notes(s, "Section turn. " + (kl.get("transition") or
                                 "State the claim, then prove it."))

    for sp in kl["supports"]:
        s = slide()
        eyebrow(s, "%s  ·  %02d" % (kl["tag"], i + 1))
        title(s, sp["assertion"], size=22)
        ex = sp.get("exhibit")
        if ex and ex["kind"] in ("bar", "column", "line"):
            txt(s, M, 1.9, 7.6, 0.3, ex.get("unit", ""), size=11, color=MUTE)
            add_chart(s, ex, M, 2.25, 7.6, 4.1)
            rect(s, M + 8.0, 2.25, CW - 8.0, 4.1, TINT, rounded=True)
            bullets(s, sp["bullets"][:3], M + 8.35, 2.6, CW - 8.7, size=12.5,
                    gap=10)
        elif ex and ex["kind"] == "table":
            add_table(s, ex, M, 2.15, 8.2)
            rect(s, M + 8.6, 2.15, CW - 8.6, 3.6, TINT, rounded=True)
            bullets(s, sp["bullets"][:3], M + 8.95, 2.5, CW - 9.3, size=12,
                    gap=10)
        else:
            for bi, b in enumerate(sp["bullets"][:3]):
                yy = 2.3 + bi * 1.28
                rect(s, M, yy, CW, 1.08, TINT, rounded=True)
                num_circle(s, M + 0.35, yy + 0.3, bi + 1, 0.46)
                txt(s, M + 1.1, yy + 0.18, CW - 1.6, 0.75, b, size=16,
                    color=INK, spacing=22)
        foot(s, "Source: " + sp["source"])
        notes(s, sp["notes"])

# ---------------- ASK ----------------
s = slide()
eyebrow(s, "So what")
title(s, "The decision requested", size=24)
rect(s, M, 1.9, CW, 1.5, SLATE, rounded=True)
txt(s, M + 0.45, 2.1, CW - 0.9, 1.1, ASK["decision"], size=20, color=WHITE,
    bold=True, font=HEAD_FONT, spacing=26)
rows = [("Gate", ASK["gate"]), ("Owner", ASK["owner"]), ("Date", ASK["date"])]
y = 3.7
for lab, text in rows:
    rect(s, M, y, CW, 0.72, TINT, rounded=True)
    txt(s, M + 0.42, y + 0.17, 1.8, 0.4, lab.upper(), size=11, color=GOLD,
        bold=True)
    txt(s, M + 2.3, y + 0.12, CW - 2.9, 0.5, text, size=13, color=INK,
        spacing=17)
    y += 0.86
txt(s, M, y + 0.1, CW, 0.5, "Key assumption: " + CONCLUSION["key_assumption"],
    size=11.5, color=MUTE, italic=True, spacing=15)
foot(s, META["footer_tag"])
notes(s, "One decision, one gate, one owner, one date. Close with the "
         "restated answer: " + CONCLUSION["restate"])

# ---------------------------------------------------------------- post-check
expected = 4 + sum(1 + len(kl["supports"]) for kl in KEY_LINE) + 1
extra = [("slide count matches the pyramid (%d expected)" % expected,
          len(prs.slides._sldIdLst) == expected),
         ("every slide carries speaker notes",
          all(sl.notes_slide.notes_text_frame.text.strip()
              for sl in prs.slides))]
if all(ok for _, ok in extra):
    prs.save(OUTPUT_FILE)
    print_report(_checks, _warnings, extra, wrote=True)
else:
    print_report(_checks, _warnings, extra, wrote=False)
